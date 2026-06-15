from __future__ import annotations

import io
import json
import re
from datetime import date, datetime, timedelta
from decimal import Decimal
from typing import Any, Optional
from zoneinfo import ZoneInfo

import httpx
from fastapi import APIRouter, Depends, Query, Response
from fastapi.responses import StreamingResponse
from pydantic import BaseModel, Field
from sqlalchemy import select, text
from sqlalchemy.ext.asyncio import AsyncSession

from config import settings
from database import get_db, ReadOnlySession
from dependencies import require_role
from services.ai_chat import national_price as national_price_svc
from services import zg_forecast_train as zg_train_svc
from models import (
    Alert,
    Bill,
    BillingCycle,
    BillingStatement,
    Category,
    ClientCanteen,
    Contract,
    DeliveryDevice,
    DeliveryVehicle,
    IoTData,
    Order,
    OrderAbnormal,
    OrderItemAllocation,
    Product,
    QualityReport,
    SupplierProductQuote,
    Tender,
    Ticket,
    User,
)
from services.ai_chat.cards import artifact_from_tool, citations_from_rag
from services.ai_chat.chat_log import log_turn
from services.ai_chat.intent import classify_intent as classify_chat_intent
from services.ai_chat.session_memory import (
    apply_memory_entities,
    load_session,
    remember_turn_summary,
    resolve_relative_day,
    seed_memory_from_history,
    update_session_from_tool,
)
from services.ai_chat.streaming import sse_events_from_result, sse_training_phase
from services.ai_chat.synthesis import should_synthesize, synthesize_answer
from services.ai_chat.tool_payload import compact_tool_result_for_llm
from services.ai_chat.rule_chain import run_rule_chain
from services.ai_chat.prompts import build_system_prompt
from services.ai_chat.rag.retriever import search_docs


router = APIRouter(prefix="/chat", tags=["ai_chat"])
_SESSION_HISTORY: dict[str, list[dict[str, str]]] = {}


class ChatMessage(BaseModel):
    role: str = "user"
    content: str = ""


class ChatRequest(BaseModel):
    session_id: Optional[str] = None
    message: Optional[str] = None
    messages: list[ChatMessage] = Field(default_factory=list)
    stream: bool = False
    # auto | national_price | system_data | how_to | report
    question_mode: str = "auto"


class ExportRequest(BaseModel):
    title: str = "监管分析报告"
    content: str = ""
    markdown: str = ""
    filename: Optional[str] = None
    format: str = "docx"
    # 用于 xlsx：结构化表格数据（来自 run_sql / 工具卡片的 columns+rows）
    columns: list[dict[str, Any]] = Field(default_factory=list)
    rows: list[dict[str, Any]] = Field(default_factory=list)


def _monitor_guard(_=Depends(require_role("monitor"))) -> None:
    return None


def _num(value: Any) -> float:
    if isinstance(value, Decimal):
        return float(value)
    try:
        return float(value or 0)
    except (TypeError, ValueError):
        return 0.0


def _today() -> date:
    return datetime.now(ZoneInfo("Asia/Shanghai")).date()


def _latest_user_text(body: ChatRequest) -> str:
    if body.message:
        return body.message.strip()
    for msg in reversed(body.messages):
        if msg.role == "user" and msg.content.strip():
            return msg.content.strip()
    return ""


def _line_rows(order: Order) -> list[dict[str, Any]]:
    items = order.items_json or []
    snaps = order.items_snapshot_json or []
    rows = []
    for idx in range(max(len(items), len(snaps))):
        item = items[idx] if idx < len(items) and isinstance(items[idx], dict) else {}
        snap = snaps[idx] if idx < len(snaps) and isinstance(snaps[idx], dict) else {}
        qty = _num(snap.get("order_quantity", item.get("quantity", item.get("qty", 0))))
        price = _num(snap.get("order_unit_price", item.get("unit_price", item.get("price", 0))))
        rows.append(
            {
                "product_id": item.get("product_id") or snap.get("product_id"),
                "goods_name": str(snap.get("product_name") or item.get("product_name") or item.get("goods_name") or f"商品#{idx + 1}"),
                "category_name": str(snap.get("category1_name") or "未归类"),
                "qty": qty,
                "amount": round(qty * price, 2),
            }
        )
    return rows


def _district_from_address(address: Optional[str]) -> str:
    text_value = address or ""
    for name in ("东城区", "西城区", "朝阳区", "丰台区", "石景山区", "海淀区", "门头沟区", "房山区", "通州区", "顺义区", "昌平区", "大兴区", "怀柔区", "平谷区", "密云区", "延庆区"):
        if name in text_value:
            return name
    return "未标注区域"


def _mask_name(value: Any) -> str:
    text_value = str(value or "").strip()
    if not text_value:
        return ""
    return text_value[0] + "**"


def _mask_phone(value: Any) -> str:
    return re.sub(r"(?<!\d)(1[3-9]\d)(\d{4})(\d{4})(?!\d)", r"\1****\3", str(value or ""))


def _repair_text(value: Any) -> str:
    text_value = str(value or "")
    if not text_value:
        return ""
    if not re.search(r"[åèéçæä]", text_value):
        return text_value
    try:
        return text_value.encode("latin1").decode("utf-8")
    except UnicodeError:
        return text_value


def _extract_ag_product_keyword(text_value: str) -> str:
    text_value = (text_value or "").strip()
    known = [
        "大白菜",
        "小白菜",
        "圆白菜",
        "奶白菜",
        "白菜",
        "山药",
        "黄瓜",
        "番茄",
        "西红柿",
        "土豆",
        "马铃薯",
        "芋头",
        "冬瓜",
        "菠菜",
        "油麦菜",
        "油菜",
        "苦菊",
        "胡萝卜",
        "红薯",
        "玉米",
        "苹果",
        "香蕉",
        "西瓜",
        "火龙果",
    ]
    for name in known:
        if name in text_value:
            return "番茄" if name == "西红柿" else "土豆" if name == "马铃薯" else name
    cleaned = re.sub(
        r"(新发地|全国农产品|中农价格|行情|价格|报价|未来|预测|趋势|今天|明天|后天|多少|怎么样|如何|查询|查一下|给我|帮我|的|近|最近|天|日|周|月|\d+)",
        "",
        text_value,
    )
    match = re.search(r"[\u4e00-\u9fff]{2,8}", cleaned)
    return match.group(0) if match else text_value


_PRICE_STOPWORDS = (
    "明天",
    "明日",
    "明儿",
    "明儿个",
    "后天",
    "未来",
    "以后",
    "下周",
    "价格",
    "价钱",
    "多少钱",
    "多少",
    "一斤",
    "一公斤",
    "什么价",
    "什么价格",
    "大概",
    "大约",
    "左右",
    "预测",
    "参考价",
    "给我",
    "请问",
    "一下",
    "怎么样",
)


def _guess_product_query(text_value: str) -> str:
    text_value = str(text_value or "").strip()
    if not text_value:
        return ""
    m = re.search(r"([\u4e00-\u9fa5A-Za-z0-9（）()·\\-\\s]{1,24})(明天|后天|未来|以后|下周)", text_value)
    candidate = m.group(1) if m else text_value
    candidate = candidate.strip("，。！？?、 ")
    for word in sorted(_PRICE_STOPWORDS, key=len, reverse=True):
        candidate = candidate.replace(word, "")
    candidate = re.sub(r"\s+", "", candidate).strip()
    if candidate in {"西红柿", "圣女果番茄"}:
        return "番茄"
    if candidate == "马铃薯":
        return "土豆"
    return candidate[:24]


def _fast_price_query_from_text(text_value: str) -> Optional[tuple[str, dict[str, Any]]]:
    text_value = str(text_value or "").strip()
    if not text_value:
        return None
    low = text_value.lower()
    price_hit = any(k in text_value for k in ("价格", "价钱", "多少钱", "什么价", "参考价")) or "price" in low
    future_hit = any(k in text_value for k in ("明天", "明日", "明儿", "后天", "未来", "以后", "下周"))
    if not (price_hit and future_hit):
        return None
    target_offset = 2 if "后天" in text_value else 1
    mode = "tomorrow" if any(k in text_value for k in ("明天", "明日", "明儿", "后天")) else "future"
    days = 7
    m = re.search(r"(未来|以后)\s*(\d{1,2})\s*天", text_value)
    if m:
        try:
            days = max(1, min(30, int(m.group(2))))
        except Exception:
            days = 7
    parsed = national_price_svc.parse_price_query(text_value)
    return "get_national_ag_forecast_price", {
        "query_text": text_value,
        "product_query": parsed.product_query,
        "mode": mode,
        "days": days,
        "target_offset": target_offset,
    }


def _fast_historical_price_from_text(text_value: str) -> Optional[tuple[str, dict[str, Any]]]:
    """昨天/前天/历史价 → 全国农产品历史行情（非预测）。"""
    t = str(text_value or "").strip()
    if not t:
        return None
    parsed = national_price_svc.parse_price_query(t)
    price_hit = any(k in t for k in ("价格", "价钱", "多少钱", "什么价", "参考价", "行情", "菜价", "批发价"))
    future_hit = any(k in t for k in ("明天", "明日", "明儿", "后天", "未来", "以后", "下周", "预测"))
    hist_hit = parsed.target_date is not None or parsed.intent in ("history", "history_range")
    if not price_hit or future_hit:
        return None
    if not (parsed.product_query or parsed.goods_name_hint):
        return None
    if parsed.product_query == t and not parsed.goods_name_hint:
        return None
    args: dict[str, Any] = {"query_text": t, "product_query": parsed.product_query}
    if parsed.target_date:
        args["target_date"] = parsed.target_date.isoformat()
    return "get_national_ag_price", args


def _extract_xinfadi_keyword(text_value: str) -> str:
    return _extract_ag_product_keyword(text_value)


def _pick_one_product(raw_query: str, names: list[str]) -> tuple[Optional[str], list[str]]:
    clean_names = [str(n).strip() for n in names if str(n).strip()]
    if not clean_names:
        return None, []
    q = str(raw_query or "").strip().lower()
    exact = [n for n in clean_names if n.lower() == q]
    if len(exact) == 1:
        return exact[0], clean_names
    starts = [n for n in clean_names if n.lower().startswith(q)]
    if len(starts) == 1:
        return starts[0], clean_names
    if len(clean_names) == 1:
        return clean_names[0], clean_names
    return None, clean_names


def _preview(obj: Any) -> str:
    try:
        value = json.dumps(obj, ensure_ascii=False)
    except Exception:
        value = str(obj)
    return value if len(value) <= 500 else value[:500] + "…"


def _fast_price_followup(text_value: str, history: list[ChatMessage], session_id: Optional[str]) -> Optional[tuple[str, dict[str, Any]]]:
    raw_history = [{"role": m.role, "content": m.content} for m in history]
    if session_id and session_id in _SESSION_HISTORY:
        raw_history = [*_SESSION_HISTORY.get(session_id, []), *raw_history]
    last_assistant = ""
    for msg in reversed(raw_history):
        if msg.get("role") == "assistant":
            last_assistant = str(msg.get("content") or "")
            break
    m = re.search(r"我找到多个可能品名：(.+?)。请你回复更准确的品名", last_assistant)
    if not m:
        return None
    candidates = [x.strip() for x in m.group(1).split("、") if x.strip()]
    current = re.sub(r"^[就查看要\s]+", "", str(text_value or "").strip())
    matched = next((c for c in candidates if c == current or c.replace(" ", "") == current.replace(" ", "")), None)
    if not matched:
        matched = next(
            (c for c in candidates if c.startswith(current) or current in c.split("（")[0].split()[0]),
            None,
        )
    if not matched:
        matched = next((c for c in candidates if current in c), None)
    if not matched:
        return None
    return "get_national_ag_forecast_price", {"query_text": matched, "mode": "future", "days": 7}


def _national_price_tool_from_text(text_value: str) -> Optional[tuple[str, dict[str, Any]]]:
    return (
        _fast_price_query_from_text(text_value)
        or _fast_historical_price_from_text(text_value)
        or _national_price_tool_from_text_slow(text_value)
    )


def _national_price_tool_from_text_slow(text_value: str) -> Optional[tuple[str, dict[str, Any]]]:
    t = str(text_value or "").strip()
    parsed = national_price_svc.parse_price_query(t)
    price_hit = any(k in t for k in ("价格", "价钱", "多少钱", "什么价", "参考价", "行情", "菜价", "批发价"))
    if not price_hit or not (parsed.product_query or parsed.goods_name_hint):
        return None
    if parsed.intent == "forecast":
        return _fast_price_query_from_text(t) or (
            "get_national_ag_forecast_price",
            {"query_text": t, "product_query": parsed.product_query, "mode": "future", "days": 7, "target_offset": 1},
        )
    args: dict[str, Any] = {"query_text": t, "product_query": parsed.product_query}
    if parsed.target_date:
        args["target_date"] = parsed.target_date.isoformat()
    return "get_national_ag_price", args


def _remember_turn(session_id: Optional[str], user_text: str, reply: str) -> None:
    if not session_id:
        return
    rows = _SESSION_HISTORY.setdefault(session_id, [])
    if user_text:
        rows.append({"role": "user", "content": user_text})
    if reply:
        rows.append({"role": "assistant", "content": reply})
    del rows[:-12]
    remember_turn_summary(session_id, user_text, reply)


def _range_from_args(args: dict[str, Any]) -> tuple[date, date]:
    today = _today()
    scope = args.get("scope")
    if scope == "today":
        return today, today
    try:
        start = date.fromisoformat(str(args.get("start_date") or ""))
    except ValueError:
        start = today - timedelta(days=29)
    try:
        end = date.fromisoformat(str(args.get("end_date") or ""))
    except ValueError:
        end = today
    if end < start:
        start, end = end, start
    return start, end


async def _orders_in_range(db: AsyncSession, start: date, end: date) -> list[Order]:
    start_dt = datetime.combine(start, datetime.min.time())
    end_dt = datetime.combine(end + timedelta(days=1), datetime.min.time())
    return (
        await db.scalars(
            select(Order).where(Order.created_at >= start_dt, Order.created_at < end_dt).order_by(Order.id.desc()).limit(5000)
        )
    ).all()


async def _names(db: AsyncSession, orders: list[Order]) -> dict[int, str]:
    user_ids = {int(o.client_id) for o in orders if o.client_id}
    users = (await db.scalars(select(User).where(User.id.in_(user_ids)))).all() if user_ids else []
    canteens = (await db.scalars(select(ClientCanteen))).all()
    out = {int(u.id): u.company_name or u.username for u in users}
    out.update({int(c.id): c.name for c in canteens})
    return out


# ============ 监管端只读工具集辅助函数 ============
# 本文件中所有 _tool_* 函数仅执行 SELECT/COUNT 查询；
# 监管端 AI 拥有最高读取权限，但绝对禁止 INSERT/UPDATE/DELETE/EXEC——任何工具违反者破坏审计可信度。


async def _user_name_map(db: AsyncSession, ids: set[int]) -> dict[int, str]:
    if not ids:
        return {}
    rows = (await db.scalars(select(User).where(User.id.in_(ids)))).all()
    return {int(u.id): (u.company_name or u.username) for u in rows}


async def _canteen_name_map(db: AsyncSession, ids: set[int]) -> dict[int, str]:
    if not ids:
        return {}
    rows = (await db.scalars(select(ClientCanteen).where(ClientCanteen.id.in_(ids)))).all()
    return {int(c.id): c.name for c in rows}


async def _product_name_map(db: AsyncSession, ids: set[int]) -> dict[int, str]:
    if not ids:
        return {}
    rows = (await db.scalars(select(Product).where(Product.id.in_(ids)))).all()
    return {int(p.id): p.name for p in rows}


async def _resolve_order_by_no_or_id(db: AsyncSession, order_no: Optional[str], order_id: Optional[int]) -> Optional[Order]:
    if order_id:
        return await db.get(Order, int(order_id))
    if order_no:
        return await db.scalar(select(Order).where(Order.order_no == order_no.strip()))
    return None


_ORDER_NO_PATTERNS = [
    re.compile(r"\b(?:OD|OR|SO)\d{6,20}\b", re.IGNORECASE),
]
_STATEMENT_NO_PATTERN = re.compile(r"\bST-?\d{6,}[A-Za-z0-9-]*\b", re.IGNORECASE)
_PRONOUN_TOKENS = (
    "这个单", "这单", "该订单", "上面那单", "刚才那单", "这张账单", "该账单", "这笔账单",
    "它", "那边", "这家", "上述", "刚才那个", "再来一份", "那个配送商", "那个学校",
)


def _extract_entities_from_history(history: list[ChatMessage], current_text: str) -> dict[str, Any]:
    """从完整 history 提取实体（order_no、statement_no），供代词与续问解析。"""
    order_nos: list[str] = []
    statement_nos: list[str] = []
    # 反向遍历，最新的实体优先
    for msg in reversed(history):
        content = (msg.content or "") if msg else ""
        if not content:
            continue
        for pat in _ORDER_NO_PATTERNS:
            for m in pat.findall(content):
                if m not in order_nos:
                    order_nos.append(m)
        for m in _STATEMENT_NO_PATTERN.findall(content):
            if m not in statement_nos:
                statement_nos.append(m)
    has_pronoun = any(tok in (current_text or "") for tok in _PRONOUN_TOKENS)
    return {
        "order_nos": order_nos[:5],
        "statement_nos": statement_nos[:5],
        "has_pronoun": has_pronoun,
        "latest_order_no": order_nos[0] if order_nos else None,
        "latest_statement_no": statement_nos[0] if statement_nos else None,
    }


async def _tool_kpi_summary(db: AsyncSession, args: dict[str, Any]) -> dict[str, Any]:
    start, end = _range_from_args(args)
    orders = await _orders_in_range(db, start, end)
    # 可选过滤：按客户/食堂/配送商/供货商切片
    client_id = args.get("client_id")
    canteen_id = args.get("canteen_id")
    delivery_id = args.get("delivery_id")
    supplier_id = args.get("supplier_id")
    if client_id is not None:
        orders = [o for o in orders if int(o.client_id or 0) == int(client_id)]
    if canteen_id is not None:
        orders = [o for o in orders if int(o.canteen_id or 0) == int(canteen_id)]
    if delivery_id is not None:
        orders = [o for o in orders if int(o.delivery_id or 0) == int(delivery_id)]
    if supplier_id is not None:
        orders = [o for o in orders if int(o.supplier_id or 0) == int(supplier_id)]
    start_dt = datetime.combine(start, datetime.min.time())
    end_dt = datetime.combine(end + timedelta(days=1), datetime.min.time())
    alerts = (
        await db.scalars(
            select(Alert).where(Alert.created_at >= start_dt, Alert.created_at < end_dt).limit(500)
        )
    ).all()
    gmv = sum(_num(o.total_amount) for o in orders)
    fulfilled = len([o for o in orders if o.status in {"收货确认", "已结算"}])
    fulfilled_rate = round(fulfilled / len(orders) * 100, 2) if orders else 0
    rows = [
        {"metric": "订单数", "value": len(orders), "unit": "单"},
        {"metric": "GMV", "value": round(gmv, 2), "unit": "元"},
        {"metric": "当日告警", "value": len(alerts), "unit": "条"},
        {"metric": "履约完成率", "value": fulfilled_rate, "unit": "%"},
    ]
    scope_hint = ""
    for k, v in (("client_id", client_id), ("canteen_id", canteen_id), ("delivery_id", delivery_id), ("supplier_id", supplier_id)):
        if v is not None:
            scope_hint += f" {k}={v}"
    return {
        "type": "kpi",
        "title": f"{start.isoformat()} 至 {end.isoformat()} 核心 KPI{scope_hint}",
        "columns": [{"key": "metric", "label": "指标"}, {"key": "value", "label": "数值"}, {"key": "unit", "label": "单位"}],
        "rows": rows,
        "summary": f"区间订单 {len(orders)} 单，GMV ¥{gmv:.2f}，当日告警 {len(alerts)} 条，履约完成率 {fulfilled_rate}%。",
    }


async def _tool_rank(db: AsyncSession, args: dict[str, Any], mode: str) -> dict[str, Any]:
    start, end = _range_from_args(args)
    limit = int(args.get("limit") or 10)
    orders = await _orders_in_range(db, start, end)
    names = await _names(db, orders)
    bucket: dict[str, dict[str, Any]] = {}
    for order in orders:
        if mode == "region":
            key = _district_from_address(order.delivery_address)
            row = bucket.setdefault(key, {"name": key, "order_count": 0, "gmv": 0.0})
            row["order_count"] += 1
            row["gmv"] += _num(order.total_amount)
        elif mode == "customer":
            key = names.get(int(order.canteen_id or order.client_id or 0), f"客户#{order.client_id}")
            row = bucket.setdefault(key, {"name": _mask_name(key), "order_count": 0, "gmv": 0.0})
            row["order_count"] += 1
            row["gmv"] += _num(order.total_amount)
        else:
            seen = set()
            for line in _line_rows(order):
                key = line["goods_name"] if mode == "goods" else line["category_name"]
                row = bucket.setdefault(key, {"name": key, "order_count": 0, "amount": 0.0, "qty": 0.0})
                row["amount"] += _num(line["amount"])
                row["qty"] += _num(line["qty"])
                seen.add(key)
            for key in seen:
                bucket[key]["order_count"] += 1
    metric = "gmv" if mode in {"region", "customer"} else "amount"
    chart_type = str(args.get("chart_type") or "bar").strip().lower()
    if chart_type not in {"pie", "bar", "line"}:
        chart_type = "bar"
    rows = sorted(bucket.values(), key=lambda row: row.get(metric, 0), reverse=True)[:limit]
    if chart_type == "pie" and len(rows) > 6:
        head, tail = rows[:5], rows[5:]
        other_row: dict[str, Any] = {"name": "其他", "order_count": 0, metric: 0.0}
        for r in tail:
            other_row["order_count"] = int(other_row.get("order_count") or 0) + int(r.get("order_count") or 0)
            other_row[metric] = _num(other_row.get(metric)) + _num(r.get(metric))
        other_row[metric] = round(_num(other_row[metric]), 2)
        rows = head + [other_row]
    for row in rows:
        for key in ("gmv", "amount", "qty"):
            if key in row:
                row[key] = round(_num(row[key]), 2)
    titles = {"region": "区域 GMV 排行", "customer": "客户 GMV 排行", "goods": "商品排行", "category": "品类排行"}
    title = titles.get(mode, "排行")
    if mode == "category" and chart_type == "pie":
        title = "一级分类销售占比"
    return {
        "type": "rank",
        "chart_type": chart_type,
        "title": title,
        "columns": [{"key": k, "label": v} for k, v in ({"name": "名称", "order_count": "订单", metric: "金额"}).items()],
        "rows": rows,
        "summary": f"已生成{title}，共 {len(rows)} 条。",
        "_pii_masked": mode == "customer",
    }


async def _tool_daily_trend(db: AsyncSession, args: dict[str, Any]) -> dict[str, Any]:
    start, end = _range_from_args(args)
    orders = await _orders_in_range(db, start, end)
    bucket: dict[str, dict[str, Any]] = {}
    day = start
    while day <= end:
        bucket[day.isoformat()] = {"date": day.isoformat(), "orders": 0, "gmv": 0.0}
        day += timedelta(days=1)
    for order in orders:
        key = order.created_at.date().isoformat()
        if key in bucket:
            bucket[key]["orders"] += 1
            bucket[key]["gmv"] += _num(order.total_amount)
    rows = list(bucket.values())
    for row in rows:
        row["gmv"] = round(row["gmv"], 2)
    return {
        "type": "trend",
        "chart_type": "line",
        "title": "GMV 趋势",
        "columns": [{"key": "date", "label": "日期"}, {"key": "orders", "label": "订单"}, {"key": "gmv", "label": "GMV"}],
        "rows": rows,
        "summary": f"已生成 {start.isoformat()} 至 {end.isoformat()} 的趋势序列。",
    }


async def _tool_ops_alerts(db: AsyncSession, args: dict[str, Any]) -> dict[str, Any]:
    limit = int(args.get("limit") or 20)
    alerts = (await db.scalars(select(Alert).order_by(Alert.id.desc()).limit(limit))).all()
    rows = [
        {
            "level": a.level,
            "type": a.type,
            "description": a.description,
            "status": a.status,
            "created_at": a.created_at.isoformat() if a.created_at else "",
        }
        for a in alerts
    ]
    return {
        "type": "alerts",
        "title": "运营预警",
        "columns": [{"key": "level", "label": "等级"}, {"key": "type", "label": "类型"}, {"key": "description", "label": "描述"}],
        "rows": rows,
        "summary": f"最近预警 {len(rows)} 条，其中开放预警 {len([r for r in rows if r['status'] == 'open'])} 条。",
    }


async def _tool_today_orders(db: AsyncSession, args: dict[str, Any]) -> dict[str, Any]:
    start, end = _today(), _today()
    orders = await _orders_in_range(db, start, end)
    user_ids = {int(o.client_id) for o in orders if o.client_id} | {int(o.delivery_id) for o in orders if o.delivery_id}
    canteen_ids = {int(o.canteen_id) for o in orders if o.canteen_id}
    user_name_map = await _user_name_map(db, user_ids)
    canteen_name_map = await _canteen_name_map(db, canteen_ids)
    rows = [
        {
            "order_id": o.id,
            "order_no": o.order_no,
            "status": o.status,
            "amount": round(_num(o.total_amount), 2),
            "client_name": user_name_map.get(int(o.client_id or 0), f"客户#{o.client_id}"),
            "canteen_name": canteen_name_map.get(int(o.canteen_id or 0), ""),
            "delivery_name": user_name_map.get(int(o.delivery_id or 0), f"配送商#{o.delivery_id}"),
            "item_count": len(o.items_json or []),
            "has_abnormal": bool(o.has_abnormal),
            "created_at": o.created_at.isoformat() if o.created_at else "",
        }
        for o in orders[:30]
    ]
    return {
        "type": "orders",
        "title": "今日订单",
        "columns": [
            {"key": "order_no", "label": "订单号"},
            {"key": "client_name", "label": "客户"},
            {"key": "canteen_name", "label": "食堂"},
            {"key": "delivery_name", "label": "配送商"},
            {"key": "status", "label": "状态"},
            {"key": "amount", "label": "金额"},
        ],
        "rows": rows,
        "summary": f"今日共 {len(orders)} 单，GMV ¥{sum(_num(o.total_amount) for o in orders):.2f}。",
    }


# ============ 类 A：实体溯源工具 ============


async def _tool_get_order_detail(db: AsyncSession, args: dict[str, Any]) -> dict[str, Any]:
    order = await _resolve_order_by_no_or_id(db, args.get("order_no"), args.get("order_id"))
    if not order:
        return {"type": "error", "rows": [], "summary": "未找到该订单，请确认订单号或 ID。"}
    user_ids = {int(x) for x in [order.client_id, order.delivery_id, order.supplier_id] if x}
    user_map = await _user_name_map(db, user_ids)
    canteen_map = await _canteen_name_map(db, {int(order.canteen_id)} if order.canteen_id else set())
    allocs = (
        await db.scalars(
            select(OrderItemAllocation).where(OrderItemAllocation.order_id == order.id)
        )
    ).all()
    supplier_ids = {int(a.supplier_id) for a in allocs if a.supplier_id}
    product_ids = {int(a.product_id) for a in allocs if a.product_id}
    supp_map = await _user_name_map(db, supplier_ids)
    prod_map = await _product_name_map(db, product_ids)
    alloc_rows = [
        {
            "line_no": a.line_no,
            "supplier_id": a.supplier_id,
            "supplier_name": supp_map.get(int(a.supplier_id or 0), f"用户#{a.supplier_id}"),
            "product_id": a.product_id,
            "product_name": prod_map.get(int(a.product_id or 0), ""),
            "quantity": _num(a.quantity),
            "unit_price": _num(a.unit_price),
            "amount": round(_num(a.quantity) * _num(a.unit_price), 2),
        }
        for a in allocs
    ]
    statements = (
        await db.scalars(
            select(BillingStatement).where(
                BillingStatement.source_snapshot_json.like(f'%"order_ids": [{order.id}%')
            )
        )
    ).all()
    statement_rows = [
        {
            "statement_no": s.statement_no,
            "role": s.role,
            "direction": s.direction,
            "status": s.status,
            "amount": _num(s.amount),
            "canteen_id": s.canteen_id,
        }
        for s in statements
    ]
    tickets = (await db.scalars(select(Ticket).where(Ticket.order_id == order.id))).all()
    ticket_rows = [
        {"id": t.id, "type": t.type, "status": t.status, "description": (t.description or "")[:100]}
        for t in tickets
    ]
    abnormals = (await db.scalars(select(OrderAbnormal).where(OrderAbnormal.order_id == order.id))).all()
    return {
        "type": "order_detail",
        "title": f"订单 {order.order_no} 详情",
        "rows": [
            {
                "order_id": order.id,
                "order_no": order.order_no,
                "status": order.status,
                "total_amount": _num(order.total_amount),
                "has_abnormal": bool(order.has_abnormal),
                "client_id": order.client_id,
                "client_name": user_map.get(int(order.client_id or 0), ""),
                "canteen_id": order.canteen_id,
                "canteen_name": canteen_map.get(int(order.canteen_id or 0), ""),
                "delivery_id": order.delivery_id,
                "delivery_name": user_map.get(int(order.delivery_id or 0), ""),
                "supplier_id": order.supplier_id,
                "supplier_name": user_map.get(int(order.supplier_id or 0), ""),
                "expected_delivery_date": order.expected_delivery_date.isoformat() if order.expected_delivery_date else None,
                "expected_delivery_slot": order.expected_delivery_slot,
                "delivery_address": order.delivery_address,
                "item_count": len(order.items_json or []),
                "created_at": order.created_at.isoformat() if order.created_at else "",
                "updated_at": order.updated_at.isoformat() if order.updated_at else "",
            }
        ],
        "allocations": alloc_rows,
        "billing_statements": statement_rows,
        "tickets": ticket_rows,
        "abnormals_count": len(abnormals),
        "summary": f"订单 {order.order_no}：客户 {user_map.get(int(order.client_id or 0),'')} / 食堂 {canteen_map.get(int(order.canteen_id or 0),'')} / 配送商 {user_map.get(int(order.delivery_id or 0),'')}，金额 ¥{_num(order.total_amount):.2f}，状态 {order.status}。分单 {len(alloc_rows)} 行；关联账单 {len(statement_rows)} 张；工单 {len(ticket_rows)} 条；异常标记 {len(abnormals)} 条。",
    }


async def _tool_get_statement_detail(db: AsyncSession, args: dict[str, Any]) -> dict[str, Any]:
    sno = (args.get("statement_no") or "").strip()
    sid = args.get("statement_id")
    row = None
    if sid:
        row = await db.get(BillingStatement, int(sid))
    elif sno:
        row = await db.scalar(select(BillingStatement).where(BillingStatement.statement_no == sno))
    if not row:
        return {"type": "error", "rows": [], "summary": "未找到该账单。"}
    user_map = await _user_name_map(db, {int(row.owner_user_id), int(row.counterparty_user_id)})
    canteen_map = await _canteen_name_map(db, {int(row.canteen_id)} if row.canteen_id else set())
    cycle = await db.get(BillingCycle, int(row.cycle_id)) if row.cycle_id else None
    snapshot = row.source_snapshot_json or {}
    order_ids = snapshot.get("order_ids") or []
    return {
        "type": "statement_detail",
        "title": f"账单 {row.statement_no} 详情",
        "rows": [
            {
                "id": row.id,
                "statement_no": row.statement_no,
                "role": row.role,
                "owner_user_id": row.owner_user_id,
                "owner_name": user_map.get(int(row.owner_user_id), ""),
                "counterparty_user_id": row.counterparty_user_id,
                "counterparty_name": user_map.get(int(row.counterparty_user_id), ""),
                "canteen_id": row.canteen_id,
                "canteen_name": canteen_map.get(int(row.canteen_id or 0), ""),
                "direction": row.direction,
                "status": row.status,
                "amount": _num(row.amount),
                "confirmed_amount": _num(row.confirmed_amount),
                "settled_amount": _num(row.settled_amount),
                "item_count": row.item_count,
                "related_order_ids": order_ids,
                "cycle_type": cycle.cycle_type if cycle else None,
                "close_day": cycle.close_day if cycle else None,
                "created_at": row.created_at.isoformat() if row.created_at else "",
                "confirmed_at": row.confirmed_at.isoformat() if row.confirmed_at else None,
            }
        ],
        "summary": f"账单 {row.statement_no}：{row.direction} ¥{_num(row.amount):.2f}，状态 {row.status}，所属食堂 {canteen_map.get(int(row.canteen_id or 0), '未设置')}，关联订单 {len(order_ids)} 张。",
    }


async def _tool_get_ticket_detail(db: AsyncSession, args: dict[str, Any]) -> dict[str, Any]:
    tid = args.get("ticket_id")
    if tid is None:
        return {"type": "error", "rows": [], "summary": "缺少 ticket_id"}
    row = await db.get(Ticket, int(tid))
    if not row:
        return {"type": "error", "rows": [], "summary": "未找到该工单。"}
    order = await db.get(Order, int(row.order_id)) if row.order_id else None
    user_map = await _user_name_map(
        db,
        {int(row.created_by)} | ({int(row.assigned_delivery_id)} if row.assigned_delivery_id else set()),
    )
    return {
        "type": "ticket_detail",
        "title": f"工单 #{row.id} 详情",
        "rows": [
            {
                "id": row.id,
                "order_id": row.order_id,
                "order_no": order.order_no if order else None,
                "type": row.type,
                "status": row.status,
                "description": row.description,
                "created_by": row.created_by,
                "created_by_name": user_map.get(int(row.created_by), ""),
                "assigned_delivery_id": row.assigned_delivery_id,
                "assigned_delivery_name": user_map.get(int(row.assigned_delivery_id or 0), ""),
                "delivery_response": row.delivery_response,
                "delivery_responded_at": row.delivery_responded_at.isoformat() if row.delivery_responded_at else None,
                "operation_resolution": row.operation_resolution,
                "operation_resolved_at": row.operation_resolved_at.isoformat() if row.operation_resolved_at else None,
                "flow_logs": row.flow_logs_json or [],
                "created_at": row.created_at.isoformat() if row.created_at else "",
            }
        ],
        "summary": f"工单 #{row.id} ({row.type}/{row.status})，关联订单 {order.order_no if order else '无'}。",
    }


async def _tool_get_contract_detail(db: AsyncSession, args: dict[str, Any]) -> dict[str, Any]:
    cid = args.get("contract_id")
    cno = (args.get("contract_no") or "").strip()
    row = None
    if cid:
        row = await db.get(Contract, int(cid))
    elif cno:
        row = await db.scalar(select(Contract).where(Contract.contract_no == cno))
    if not row:
        return {"type": "error", "rows": [], "summary": "未找到该合约。"}
    user_map = await _user_name_map(db, {int(row.client_id), int(row.delivery_id)})
    canteens = (
        await db.scalars(select(ClientCanteen).where(ClientCanteen.school_client_id == row.client_id))
    ).all()
    canteen_rows = [{"id": c.id, "name": c.name, "status": c.status} for c in canteens]
    cat_ids = row.category_ids_json or []
    cats = (await db.scalars(select(Category).where(Category.id.in_(cat_ids)))).all() if cat_ids else []
    return {
        "type": "contract_detail",
        "title": f"合约 {row.contract_no} 详情",
        "rows": [
            {
                "id": row.id,
                "contract_no": row.contract_no,
                "client_id": row.client_id,
                "client_name": user_map.get(int(row.client_id), ""),
                "delivery_id": row.delivery_id,
                "delivery_name": user_map.get(int(row.delivery_id), ""),
                "status": row.status,
                "period_start": row.period_start.isoformat() if row.period_start else None,
                "period_end": row.period_end.isoformat() if row.period_end else None,
                "price_float_rate": row.price_float_rate,
                "category_count": len(cat_ids),
                "category_names": [c.name for c in cats],
                "category_rates_json": row.category_rates_json,
            }
        ],
        "canteens": canteen_rows,
        "summary": f"合约 {row.contract_no}：{user_map.get(int(row.client_id), '')} ↔ {user_map.get(int(row.delivery_id), '')}，状态 {row.status}，{row.period_start}~{row.period_end}，涉及 {len(canteen_rows)} 个食堂。",
    }


async def _tool_get_user_detail(db: AsyncSession, args: dict[str, Any]) -> dict[str, Any]:
    uid = args.get("user_id")
    uname = (args.get("username") or "").strip()
    row = None
    if uid:
        row = await db.get(User, int(uid))
    elif uname:
        row = await db.scalar(select(User).where(User.username == uname))
    if not row:
        return {"type": "error", "rows": [], "summary": "未找到该用户。"}
    extra: dict[str, Any] = {}
    if row.role == "client":
        canteens = (
            await db.scalars(select(ClientCanteen).where(ClientCanteen.school_client_id == row.id))
        ).all()
        extra["canteens"] = [{"id": c.id, "name": c.name, "status": c.status} for c in canteens]
    elif row.role == "delivery":
        vehicles = (
            await db.scalars(select(DeliveryVehicle).where(DeliveryVehicle.delivery_id == row.id))
        ).all()
        devices = (
            await db.scalars(select(DeliveryDevice).where(DeliveryDevice.delivery_id == row.id))
        ).all()
        bound_suppliers = (
            await db.scalars(select(User).where(User.role == "supplier", User.supplier_delivery_id == row.id))
        ).all()
        extra["vehicle_count"] = len(vehicles)
        extra["device_count"] = len(devices)
        extra["bound_suppliers"] = [
            {"id": s.id, "name": s.company_name or s.username} for s in bound_suppliers
        ]
    elif row.role == "supplier":
        quote_cnt = await db.scalar(
            select(func_count_id()).select_from(SupplierProductQuote).where(
                SupplierProductQuote.supplier_id == row.id
            )
        )
        extra["quote_count"] = int(quote_cnt or 0)
        extra["bound_delivery_id"] = row.supplier_delivery_id
    return {
        "type": "user_detail",
        "title": f"用户 {row.username} 详情",
        "rows": [
            {
                "id": row.id,
                "username": row.username,
                "role": row.role,
                "company_name": row.company_name,
                "address": row.address,
                "status": row.status,
                "contact_phone": _mask_phone(row.contact_phone) if row.contact_phone else "",
                "created_at": row.created_at.isoformat() if getattr(row, "created_at", None) else "",
            }
        ],
        **extra,
        "summary": f"用户 {row.username} ({row.role})，公司 {row.company_name or '-'}，状态 {row.status}。",
    }


# 工具内部辅助：避免重复导入 func.count
def func_count_id():
    from sqlalchemy import func as _func
    return _func.count()


# ============ 类 B：列表搜索工具 ============


def _clamp_limit(args: dict[str, Any], default: int = 30, cap: int = 50) -> int:
    raw = args.get("limit") or default
    try:
        return max(1, min(cap, int(raw)))
    except (TypeError, ValueError):
        return default


def _parse_date(value: Any) -> Optional[date]:
    if not value:
        return None
    if isinstance(value, date):
        return value
    try:
        return date.fromisoformat(str(value).strip()[:10])
    except (TypeError, ValueError):
        return None


async def _tool_search_orders(db: AsyncSession, args: dict[str, Any]) -> dict[str, Any]:
    limit = _clamp_limit(args, default=20, cap=50)
    start = _parse_date(args.get("start_date")) or (_today() - timedelta(days=7))
    end = _parse_date(args.get("end_date")) or _today()
    if end < start:
        start, end = end, start
    start_dt = datetime.combine(start, datetime.min.time())
    end_dt = datetime.combine(end + timedelta(days=1), datetime.min.time())
    stmt = select(Order).where(Order.created_at >= start_dt, Order.created_at < end_dt)
    if args.get("client_id") is not None:
        stmt = stmt.where(Order.client_id == int(args["client_id"]))
    if args.get("canteen_id") is not None:
        stmt = stmt.where(Order.canteen_id == int(args["canteen_id"]))
    if args.get("delivery_id") is not None:
        stmt = stmt.where(Order.delivery_id == int(args["delivery_id"]))
    if args.get("supplier_id") is not None:
        stmt = stmt.where(Order.supplier_id == int(args["supplier_id"]))
    if args.get("status"):
        stmt = stmt.where(Order.status == str(args["status"]))
    if args.get("has_abnormal") is True:
        stmt = stmt.where(Order.has_abnormal.is_(True))
    if args.get("amount_min") is not None:
        stmt = stmt.where(Order.total_amount >= Decimal(str(args["amount_min"])))
    if args.get("amount_max") is not None:
        stmt = stmt.where(Order.total_amount <= Decimal(str(args["amount_max"])))
    if args.get("order_no_like"):
        stmt = stmt.where(Order.order_no.like(f"%{args['order_no_like']}%"))
    rows = (await db.scalars(stmt.order_by(Order.id.desc()).limit(limit))).all()
    user_ids = {int(o.client_id) for o in rows if o.client_id} | {int(o.delivery_id) for o in rows if o.delivery_id}
    canteen_ids = {int(o.canteen_id) for o in rows if o.canteen_id}
    user_map = await _user_name_map(db, user_ids)
    canteen_map = await _canteen_name_map(db, canteen_ids)
    out = [
        {
            "order_id": o.id,
            "order_no": o.order_no,
            "status": o.status,
            "amount": round(_num(o.total_amount), 2),
            "client_name": user_map.get(int(o.client_id or 0), ""),
            "canteen_name": canteen_map.get(int(o.canteen_id or 0), ""),
            "delivery_name": user_map.get(int(o.delivery_id or 0), ""),
            "has_abnormal": bool(o.has_abnormal),
            "created_at": o.created_at.isoformat() if o.created_at else "",
        }
        for o in rows
    ]
    return {
        "type": "orders",
        "title": "订单列表",
        "columns": [
            {"key": "order_no", "label": "订单号"},
            {"key": "client_name", "label": "客户"},
            {"key": "canteen_name", "label": "食堂"},
            {"key": "delivery_name", "label": "配送商"},
            {"key": "status", "label": "状态"},
            {"key": "amount", "label": "金额"},
        ],
        "rows": out,
        "summary": f"匹配到 {len(out)} 单（{start} ~ {end}），GMV ¥{sum(r['amount'] for r in out):.2f}。",
    }


async def _tool_search_statements(db: AsyncSession, args: dict[str, Any]) -> dict[str, Any]:
    limit = _clamp_limit(args, default=20, cap=50)
    stmt = select(BillingStatement)
    if args.get("direction"):
        stmt = stmt.where(BillingStatement.direction == str(args["direction"]))
    if args.get("status"):
        stmt = stmt.where(BillingStatement.status == str(args["status"]))
    if args.get("owner_user_id") is not None:
        stmt = stmt.where(BillingStatement.owner_user_id == int(args["owner_user_id"]))
    if args.get("counterparty_user_id") is not None:
        stmt = stmt.where(BillingStatement.counterparty_user_id == int(args["counterparty_user_id"]))
    if args.get("canteen_id") is not None:
        stmt = stmt.where(BillingStatement.canteen_id == int(args["canteen_id"]))
    if args.get("role"):
        stmt = stmt.where(BillingStatement.role == str(args["role"]))
    start = _parse_date(args.get("start_date"))
    end = _parse_date(args.get("end_date"))
    if start:
        stmt = stmt.where(BillingStatement.created_at >= datetime.combine(start, datetime.min.time()))
    if end:
        stmt = stmt.where(BillingStatement.created_at < datetime.combine(end + timedelta(days=1), datetime.min.time()))
    rows = (await db.scalars(stmt.order_by(BillingStatement.id.desc()).limit(limit))).all()
    user_ids = {int(r.owner_user_id) for r in rows} | {int(r.counterparty_user_id) for r in rows}
    canteen_ids = {int(r.canteen_id) for r in rows if r.canteen_id}
    user_map = await _user_name_map(db, user_ids)
    canteen_map = await _canteen_name_map(db, canteen_ids)
    out = [
        {
            "id": r.id,
            "statement_no": r.statement_no,
            "role": r.role,
            "direction": r.direction,
            "status": r.status,
            "amount": _num(r.amount),
            "settled_amount": _num(r.settled_amount),
            "owner_name": user_map.get(int(r.owner_user_id), ""),
            "counterparty_name": user_map.get(int(r.counterparty_user_id), ""),
            "canteen_name": canteen_map.get(int(r.canteen_id or 0), ""),
            "created_at": r.created_at.isoformat() if r.created_at else "",
        }
        for r in rows
    ]
    return {
        "type": "statements",
        "title": "账单列表",
        "rows": out,
        "summary": f"匹配账单 {len(out)} 张，合计金额 ¥{sum(_num(r['amount']) for r in out):.2f}，未结清 ¥{sum(max(0, _num(r['amount']) - _num(r['settled_amount'])) for r in out):.2f}。",
    }


async def _tool_search_tickets(db: AsyncSession, args: dict[str, Any]) -> dict[str, Any]:
    limit = _clamp_limit(args, default=20, cap=50)
    stmt = select(Ticket)
    if args.get("status"):
        stmt = stmt.where(Ticket.status == str(args["status"]))
    if args.get("type"):
        stmt = stmt.where(Ticket.type == str(args["type"]))
    if args.get("order_id") is not None:
        stmt = stmt.where(Ticket.order_id == int(args["order_id"]))
    start = _parse_date(args.get("start_date"))
    end = _parse_date(args.get("end_date"))
    if start:
        stmt = stmt.where(Ticket.created_at >= datetime.combine(start, datetime.min.time()))
    if end:
        stmt = stmt.where(Ticket.created_at < datetime.combine(end + timedelta(days=1), datetime.min.time()))
    rows = (await db.scalars(stmt.order_by(Ticket.id.desc()).limit(limit))).all()
    order_ids = {int(t.order_id) for t in rows if t.order_id}
    orders = (await db.scalars(select(Order).where(Order.id.in_(order_ids)))).all() if order_ids else []
    order_no_map = {int(o.id): o.order_no for o in orders}
    out = [
        {
            "id": t.id,
            "order_id": t.order_id,
            "order_no": order_no_map.get(int(t.order_id or 0), ""),
            "type": t.type,
            "status": t.status,
            "description": (t.description or "")[:120],
            "assigned_delivery_id": t.assigned_delivery_id,
            "created_at": t.created_at.isoformat() if t.created_at else "",
        }
        for t in rows
    ]
    return {
        "type": "tickets",
        "title": "工单列表",
        "rows": out,
        "summary": f"匹配工单 {len(out)} 条，其中待处理 {sum(1 for r in out if r['status']=='待处理')} 条，处理中 {sum(1 for r in out if r['status']=='处理中')} 条。",
    }


async def _tool_search_quality_reports(db: AsyncSession, args: dict[str, Any]) -> dict[str, Any]:
    limit = _clamp_limit(args, default=20, cap=50)
    stmt = select(QualityReport)
    if args.get("supplier_id") is not None:
        stmt = stmt.where(QualityReport.supplier_id == int(args["supplier_id"]))
    if args.get("order_id") is not None:
        stmt = stmt.where(QualityReport.order_id == int(args["order_id"]))
    if args.get("product_id") is not None:
        stmt = stmt.where(QualityReport.product_id == int(args["product_id"]))
    if args.get("status"):
        stmt = stmt.where(QualityReport.status == str(args["status"]))
    start = _parse_date(args.get("start_date"))
    end = _parse_date(args.get("end_date"))
    if start:
        stmt = stmt.where(QualityReport.created_at >= datetime.combine(start, datetime.min.time()))
    if end:
        stmt = stmt.where(QualityReport.created_at < datetime.combine(end + timedelta(days=1), datetime.min.time()))
    rows = (await db.scalars(stmt.order_by(QualityReport.id.desc()).limit(limit))).all()
    user_ids = {int(r.supplier_id) for r in rows if r.supplier_id}
    product_ids = {int(r.product_id) for r in rows if r.product_id}
    user_map = await _user_name_map(db, user_ids)
    prod_map = await _product_name_map(db, product_ids)
    out = [
        {
            "id": r.id,
            "report_no": r.report_no,
            "supplier_id": r.supplier_id,
            "supplier_name": user_map.get(int(r.supplier_id or 0), ""),
            "product_id": r.product_id,
            "product_name": prod_map.get(int(r.product_id or 0), ""),
            "order_id": r.order_id,
            "status": r.status,
            "file_url": r.file_url,
            "created_at": r.created_at.isoformat() if r.created_at else "",
        }
        for r in rows
    ]
    return {
        "type": "quality_reports",
        "title": "质检报告列表",
        "rows": out,
        "summary": f"匹配质检报告 {len(out)} 条，待审核 {sum(1 for r in out if r['status']=='待审核')} 条。",
    }


async def _tool_search_alerts(db: AsyncSession, args: dict[str, Any]) -> dict[str, Any]:
    limit = _clamp_limit(args, default=20, cap=50)
    stmt = select(Alert)
    if args.get("level"):
        stmt = stmt.where(Alert.level == str(args["level"]))
    if args.get("status"):
        stmt = stmt.where(Alert.status == str(args["status"]))
    if args.get("type"):
        stmt = stmt.where(Alert.type == str(args["type"]))
    start = _parse_date(args.get("start_date"))
    end = _parse_date(args.get("end_date"))
    if start:
        stmt = stmt.where(Alert.created_at >= datetime.combine(start, datetime.min.time()))
    if end:
        stmt = stmt.where(Alert.created_at < datetime.combine(end + timedelta(days=1), datetime.min.time()))
    rows = (await db.scalars(stmt.order_by(Alert.id.desc()).limit(limit))).all()
    out = [
        {
            "id": r.id,
            "level": r.level,
            "type": r.type,
            "description": (r.description or "")[:150],
            "status": r.status,
            "created_at": r.created_at.isoformat() if r.created_at else "",
        }
        for r in rows
    ]
    return {
        "type": "alerts",
        "title": "告警列表",
        "rows": out,
        "summary": f"匹配告警 {len(out)} 条，开放 {sum(1 for r in out if r['status']=='open')} 条，高危 {sum(1 for r in out if r['level']=='high')} 条。",
    }


# ============ 类 C：主体聚合指标 ============


def _default_metric_range(args: dict[str, Any]) -> tuple[date, date]:
    start = _parse_date(args.get("start_date")) or (_today() - timedelta(days=29))
    end = _parse_date(args.get("end_date")) or _today()
    if end < start:
        start, end = end, start
    return start, end


async def _orders_filtered_in_range(
    db: AsyncSession, start: date, end: date, **filters
) -> list[Order]:
    start_dt = datetime.combine(start, datetime.min.time())
    end_dt = datetime.combine(end + timedelta(days=1), datetime.min.time())
    stmt = select(Order).where(Order.created_at >= start_dt, Order.created_at < end_dt)
    for col, val in filters.items():
        if val is not None:
            stmt = stmt.where(getattr(Order, col) == val)
    return (await db.scalars(stmt.order_by(Order.id.desc()).limit(5000))).all()


async def _tool_get_delivery_metrics(db: AsyncSession, args: dict[str, Any]) -> dict[str, Any]:
    did = args.get("delivery_id")
    if did is None:
        return {"type": "error", "rows": [], "summary": "缺少 delivery_id"}
    start, end = _default_metric_range(args)
    orders = await _orders_filtered_in_range(db, start, end, delivery_id=int(did))
    gmv = sum(_num(o.total_amount) for o in orders)
    fulfilled = [o for o in orders if o.status in {"收货确认", "已结算"}]
    abnormal = [o for o in orders if o.has_abnormal]
    on_time = 0
    for o in fulfilled:
        if o.expected_delivery_date and o.updated_at:
            if o.updated_at.date() <= o.expected_delivery_date:
                on_time += 1
    # 客户分布 top5
    client_ids = {int(o.client_id) for o in orders if o.client_id}
    user_map = await _user_name_map(db, client_ids)
    client_bucket: dict[int, dict[str, Any]] = {}
    for o in orders:
        cid = int(o.client_id or 0)
        b = client_bucket.setdefault(cid, {"client_id": cid, "client_name": user_map.get(cid, ""), "order_count": 0, "gmv": 0.0})
        b["order_count"] += 1
        b["gmv"] += _num(o.total_amount)
    top_clients = sorted(client_bucket.values(), key=lambda x: x["gmv"], reverse=True)[:5]
    for c in top_clients:
        c["gmv"] = round(c["gmv"], 2)
    return {
        "type": "metrics",
        "title": f"配送商履约指标 ({start} ~ {end})",
        "rows": [
            {"metric": "履约单数", "value": len(orders), "unit": "单"},
            {"metric": "GMV", "value": round(gmv, 2), "unit": "元"},
            {"metric": "完成单数", "value": len(fulfilled), "unit": "单"},
            {"metric": "完成率", "value": round(len(fulfilled) / len(orders) * 100, 2) if orders else 0, "unit": "%"},
            {"metric": "按期率", "value": round(on_time / len(fulfilled) * 100, 2) if fulfilled else 0, "unit": "%"},
            {"metric": "异常率", "value": round(len(abnormal) / len(orders) * 100, 2) if orders else 0, "unit": "%"},
        ],
        "top_clients": top_clients,
        "summary": f"配送商在 {start}~{end} 履约 {len(orders)} 单，GMV ¥{gmv:.2f}，完成率 {round(len(fulfilled)/len(orders)*100, 2) if orders else 0}%，异常 {len(abnormal)} 单。",
    }


async def _tool_get_supplier_metrics(db: AsyncSession, args: dict[str, Any]) -> dict[str, Any]:
    sid = args.get("supplier_id")
    if sid is None:
        return {"type": "error", "rows": [], "summary": "缺少 supplier_id"}
    start, end = _default_metric_range(args)
    start_dt = datetime.combine(start, datetime.min.time())
    end_dt = datetime.combine(end + timedelta(days=1), datetime.min.time())
    # 通过 OrderItemAllocation join Order 找该 supplier 涉及的订单
    allocs = (
        await db.scalars(
            select(OrderItemAllocation)
            .join(Order, Order.id == OrderItemAllocation.order_id)
            .where(
                OrderItemAllocation.supplier_id == int(sid),
                Order.created_at >= start_dt,
                Order.created_at < end_dt,
            )
            .limit(5000)
        )
    ).all()
    order_ids = {int(a.order_id) for a in allocs}
    total_amount = sum(_num(a.quantity) * _num(a.unit_price) for a in allocs)
    prod_ids = {int(a.product_id) for a in allocs}
    prod_map = await _product_name_map(db, prod_ids)
    prod_bucket: dict[int, dict[str, Any]] = {}
    for a in allocs:
        pid = int(a.product_id or 0)
        b = prod_bucket.setdefault(pid, {"product_id": pid, "product_name": prod_map.get(pid, ""), "qty": 0.0, "amount": 0.0})
        b["qty"] += _num(a.quantity)
        b["amount"] += _num(a.quantity) * _num(a.unit_price)
    top_products = sorted(prod_bucket.values(), key=lambda x: x["amount"], reverse=True)[:10]
    for p in top_products:
        p["qty"] = round(p["qty"], 2)
        p["amount"] = round(p["amount"], 2)
    return {
        "type": "metrics",
        "title": f"供货商指标 ({start} ~ {end})",
        "rows": [
            {"metric": "涉及订单", "value": len(order_ids), "unit": "单"},
            {"metric": "分单行数", "value": len(allocs), "unit": "行"},
            {"metric": "累计金额", "value": round(total_amount, 2), "unit": "元"},
            {"metric": "商品种类", "value": len(prod_ids), "unit": "种"},
        ],
        "top_products": top_products,
        "summary": f"供货商在 {start}~{end} 涉及 {len(order_ids)} 单，分单 {len(allocs)} 行，累计 ¥{total_amount:.2f}，商品 {len(prod_ids)} 种。",
    }


async def _tool_get_client_metrics(db: AsyncSession, args: dict[str, Any]) -> dict[str, Any]:
    cid = args.get("client_id")
    if cid is None:
        return {"type": "error", "rows": [], "summary": "缺少 client_id"}
    start, end = _default_metric_range(args)
    filters: dict[str, Any] = {"client_id": int(cid)}
    if args.get("canteen_id") is not None:
        filters["canteen_id"] = int(args["canteen_id"])
    orders = await _orders_filtered_in_range(db, start, end, **filters)
    gmv = sum(_num(o.total_amount) for o in orders)
    # 按食堂拆分
    canteen_ids = {int(o.canteen_id) for o in orders if o.canteen_id}
    canteen_map = await _canteen_name_map(db, canteen_ids)
    canteen_bucket: dict[int, dict[str, Any]] = {}
    for o in orders:
        canid = int(o.canteen_id or 0)
        b = canteen_bucket.setdefault(canid, {"canteen_id": canid, "canteen_name": canteen_map.get(canid, ""), "order_count": 0, "gmv": 0.0})
        b["order_count"] += 1
        b["gmv"] += _num(o.total_amount)
    by_canteen = sorted(canteen_bucket.values(), key=lambda x: x["gmv"], reverse=True)
    for r in by_canteen:
        r["gmv"] = round(r["gmv"], 2)
    # 按品类拆分
    cat_bucket: dict[str, dict[str, Any]] = {}
    for o in orders:
        for line in _line_rows(o):
            key = line.get("category_name") or "未分类"
            b = cat_bucket.setdefault(key, {"category_name": key, "amount": 0.0, "qty": 0.0})
            b["amount"] += _num(line.get("amount"))
            b["qty"] += _num(line.get("qty"))
    by_category = sorted(cat_bucket.values(), key=lambda x: x["amount"], reverse=True)[:10]
    for r in by_category:
        r["amount"] = round(r["amount"], 2)
        r["qty"] = round(r["qty"], 2)
    # 配送商分布
    delivery_ids = {int(o.delivery_id) for o in orders if o.delivery_id}
    user_map = await _user_name_map(db, delivery_ids)
    delivery_bucket: dict[int, dict[str, Any]] = {}
    for o in orders:
        did = int(o.delivery_id or 0)
        b = delivery_bucket.setdefault(did, {"delivery_id": did, "delivery_name": user_map.get(did, ""), "order_count": 0, "gmv": 0.0})
        b["order_count"] += 1
        b["gmv"] += _num(o.total_amount)
    by_delivery = sorted(delivery_bucket.values(), key=lambda x: x["gmv"], reverse=True)
    for r in by_delivery:
        r["gmv"] = round(r["gmv"], 2)
    return {
        "type": "metrics",
        "title": f"客户采购指标 ({start} ~ {end})",
        "rows": [
            {"metric": "订单数", "value": len(orders), "unit": "单"},
            {"metric": "GMV", "value": round(gmv, 2), "unit": "元"},
            {"metric": "涉及食堂", "value": len(canteen_bucket), "unit": "个"},
        ],
        "by_canteen": by_canteen,
        "by_category": by_category,
        "by_delivery": by_delivery,
        "summary": f"客户在 {start}~{end} 采购 {len(orders)} 单 ¥{gmv:.2f}，涉及 {len(canteen_bucket)} 个食堂、{len(delivery_bucket)} 个配送商。",
    }


async def _tool_get_factory_metrics(db: AsyncSession, args: dict[str, Any]) -> dict[str, Any]:
    fid = args.get("factory_id")
    if fid is None:
        return {"type": "error", "rows": [], "summary": "缺少 factory_id"}
    start, end = _default_metric_range(args)
    start_dt = datetime.combine(start, datetime.min.time())
    end_dt = datetime.combine(end + timedelta(days=1), datetime.min.time())
    # 厂家也是 supplier_id（OrderItemAllocation.supplier_id 字段，厂家时该 id 是 factory user id）
    allocs = (
        await db.scalars(
            select(OrderItemAllocation)
            .join(Order, Order.id == OrderItemAllocation.order_id)
            .where(
                OrderItemAllocation.supplier_id == int(fid),
                Order.created_at >= start_dt,
                Order.created_at < end_dt,
            )
            .limit(5000)
        )
    ).all()
    order_ids = {int(a.order_id) for a in allocs}
    total_amount = sum(_num(a.quantity) * _num(a.unit_price) for a in allocs)
    return {
        "type": "metrics",
        "title": f"厂家分单指标 ({start} ~ {end})",
        "rows": [
            {"metric": "涉及订单", "value": len(order_ids), "unit": "单"},
            {"metric": "分单行数", "value": len(allocs), "unit": "行"},
            {"metric": "累计金额", "value": round(total_amount, 2), "unit": "元"},
        ],
        "summary": f"厂家在 {start}~{end} 涉及 {len(order_ids)} 单，分单 {len(allocs)} 行，累计 ¥{total_amount:.2f}。",
    }


# ============ 类 D：IoT/资产/账务监督 ============


async def _tool_get_device_status(db: AsyncSession, args: dict[str, Any]) -> dict[str, Any]:
    limit = _clamp_limit(args, default=30, cap=50)
    stmt = select(DeliveryDevice)
    if args.get("device_type"):
        stmt = stmt.where(DeliveryDevice.device_type == str(args["device_type"]))
    if args.get("delivery_id") is not None:
        stmt = stmt.where(DeliveryDevice.delivery_id == int(args["delivery_id"]))
    if args.get("status"):
        stmt = stmt.where(DeliveryDevice.status == str(args["status"]))
    rows = (await db.scalars(stmt.order_by(DeliveryDevice.id.desc()).limit(limit))).all()
    # 取每个设备最近一条 IoT 数据时间作为心跳（按 device_code 关联 iot_data.device_id）
    last_seen: dict[str, datetime] = {}
    for r in rows:
        last = await db.scalar(
            select(IoTData.recorded_at)
            .where(IoTData.device_id == r.device_code)
            .order_by(IoTData.id.desc())
            .limit(1)
        )
        if last:
            last_seen[r.device_code] = last
    now = datetime.utcnow()
    out = []
    for r in rows:
        last = last_seen.get(r.device_code)
        offline_minutes = int((now - last).total_seconds() // 60) if last else None
        out.append({
            "id": r.id,
            "delivery_id": r.delivery_id,
            "device_type": r.device_type,
            "vendor": r.vendor,
            "device_code": r.device_code,
            "device_name": r.device_name,
            "status": r.status,
            "last_seen_at": last.isoformat() if last else None,
            "offline_minutes": offline_minutes,
        })
    if args.get("offline_only"):
        out = [r for r in out if r["offline_minutes"] is None or r["offline_minutes"] > 10]
    if args.get("online_only"):
        out = [r for r in out if r["offline_minutes"] is not None and r["offline_minutes"] <= 10]
    return {
        "type": "devices",
        "title": "设备状态",
        "rows": out,
        "summary": f"返回设备 {len(out)} 台，离线/未上报 {sum(1 for r in out if r['offline_minutes'] is None or r['offline_minutes'] > 10)} 台。",
    }


async def _tool_get_iot_recent(db: AsyncSession, args: dict[str, Any]) -> dict[str, Any]:
    limit = _clamp_limit(args, default=30, cap=50)
    hours = args.get("hours") or 1
    try:
        hours = max(1, min(168, int(hours)))
    except (TypeError, ValueError):
        hours = 1
    since = datetime.utcnow() - timedelta(hours=hours)
    stmt = select(IoTData).where(IoTData.recorded_at >= since)
    if args.get("device_type"):
        stmt = stmt.where(IoTData.device_type == str(args["device_type"]))
    if args.get("device_id"):
        stmt = stmt.where(IoTData.device_id == str(args["device_id"]))
    rows = (await db.scalars(stmt.order_by(IoTData.id.desc()).limit(limit))).all()
    out = [
        {
            "id": r.id,
            "device_type": r.device_type,
            "device_id": r.device_id,
            "recorded_at": r.recorded_at.isoformat() if r.recorded_at else "",
            "payload": r.payload_json,
        }
        for r in rows
    ]
    return {
        "type": "iot",
        "title": f"近 {hours} 小时 IoT 数据",
        "rows": out,
        "summary": f"返回 {len(out)} 条记录（近 {hours}h）。",
    }


async def _tool_get_overdue_bills(db: AsyncSession, args: dict[str, Any]) -> dict[str, Any]:
    from services.billing_service import compute_billing_period, summarize_overdue

    scope = (args.get("scope") or "payment").strip().lower()
    limit = _clamp_limit(args, default=20, cap=50)
    rows = (
        await db.scalars(select(BillingStatement).order_by(BillingStatement.id.desc()).limit(500))
    ).all()
    cycle_ids = {int(r.cycle_id) for r in rows if r.cycle_id}
    cycles = (await db.scalars(select(BillingCycle).where(BillingCycle.id.in_(cycle_ids)))).all() if cycle_ids else []
    cycle_map = {int(c.id): c for c in cycles}
    enriched: list[dict[str, Any]] = []
    for r in rows:
        cycle = cycle_map.get(int(r.cycle_id or 0))
        if not cycle:
            continue
        period = compute_billing_period(r.created_at or datetime.utcnow(), cycle)
        enriched.append({
            "id": r.id,
            "statement_no": r.statement_no,
            "role": r.role,
            "direction": r.direction,
            "status": r.status,
            "amount": _num(r.amount),
            "settled_amount": _num(r.settled_amount),
            "canteen_id": r.canteen_id,
            "owner_user_id": r.owner_user_id,
            "counterparty_user_id": r.counterparty_user_id,
            **period,
        })
    today = _today()
    overdue = summarize_overdue(enriched, today)
    bucket = overdue["overdue_payment"] if scope == "payment" else overdue["overdue_confirm"]
    bucket = bucket[:limit]
    user_ids = {int(r["owner_user_id"]) for r in bucket} | {int(r["counterparty_user_id"]) for r in bucket}
    user_map = await _user_name_map(db, user_ids)
    canteen_ids = {int(r["canteen_id"]) for r in bucket if r.get("canteen_id")}
    canteen_map = await _canteen_name_map(db, canteen_ids)
    for r in bucket:
        r["owner_name"] = user_map.get(int(r["owner_user_id"]), "")
        r["counterparty_name"] = user_map.get(int(r["counterparty_user_id"]), "")
        r["canteen_name"] = canteen_map.get(int(r.get("canteen_id") or 0), "")
    return {
        "type": "statements",
        "title": ("超期未付款" if scope == "payment" else "超期未确认") + " 账单",
        "rows": bucket,
        "summary": f"超期{'未付' if scope=='payment' else '未确认'} {len(bucket)} 张，合计 ¥{sum(_num(r['amount']) for r in bucket):.2f}。",
    }


# ============ 类 E：报价/合约对比 ============


async def _tool_get_supplier_quotes(db: AsyncSession, args: dict[str, Any]) -> dict[str, Any]:
    limit = _clamp_limit(args, default=30, cap=50)
    stmt = select(SupplierProductQuote)
    if args.get("supplier_id") is not None:
        stmt = stmt.where(SupplierProductQuote.supplier_id == int(args["supplier_id"]))
    if args.get("product_id") is not None:
        stmt = stmt.where(SupplierProductQuote.product_id == int(args["product_id"]))
    rows = (await db.scalars(stmt.order_by(SupplierProductQuote.id.desc()).limit(limit))).all()
    sup_ids = {int(r.supplier_id) for r in rows}
    prod_ids = {int(r.product_id) for r in rows}
    user_map = await _user_name_map(db, sup_ids)
    prod_map = await _product_name_map(db, prod_ids)
    out = [
        {
            "id": r.id,
            "supplier_id": r.supplier_id,
            "supplier_name": user_map.get(int(r.supplier_id), ""),
            "product_id": r.product_id,
            "product_name": prod_map.get(int(r.product_id), ""),
            "quote_price": _num(r.quote_price),
            "remark": r.remark,
            "created_at": r.created_at.isoformat() if r.created_at else "",
        }
        for r in rows
    ]
    return {
        "type": "quotes",
        "title": "供货商报价列表",
        "rows": out,
        "summary": f"匹配报价 {len(out)} 条。",
    }


async def _tool_compare_supplier_vs_xinfadi(db: AsyncSession, args: dict[str, Any]) -> dict[str, Any]:
    pid = args.get("product_id")
    if pid is None:
        return {"type": "error", "rows": [], "summary": "缺少 product_id"}
    product = await db.get(Product, int(pid))
    if not product:
        return {"type": "error", "rows": [], "summary": "未找到该商品。"}
    stmt = select(SupplierProductQuote).where(SupplierProductQuote.product_id == int(pid))
    if args.get("supplier_id") is not None:
        stmt = stmt.where(SupplierProductQuote.supplier_id == int(args["supplier_id"]))
    quotes = (await db.scalars(stmt.limit(50))).all()
    sup_map = await _user_name_map(db, {int(q.supplier_id) for q in quotes})
    hist = await national_price_svc.history_for_query(db, product.name)
    hist_rows = hist.get("rows") or []
    avg_price = round(sum(_num(r.get("avg_price")) for r in hist_rows) / len(hist_rows), 2) if hist_rows else None
    rows = [
        {
            "supplier_id": q.supplier_id,
            "supplier_name": sup_map.get(int(q.supplier_id), ""),
            "quote_price": _num(q.quote_price),
            "reference_price": _num(product.reference_price),
            "xinfadi_avg_price": avg_price,
            "vs_reference_pct": round((_num(q.quote_price) - _num(product.reference_price)) / _num(product.reference_price) * 100, 2) if _num(product.reference_price) else None,
            "vs_xinfadi_pct": round((_num(q.quote_price) - avg_price) / avg_price * 100, 2) if avg_price else None,
        }
        for q in quotes
    ]
    return {
        "type": "quotes_compare",
        "title": f"{product.name} 报价对比",
        "rows": rows,
        "reference": {"product_name": product.name, "reference_price": _num(product.reference_price), "xinfadi_avg_price": avg_price},
        "summary": f"商品 {product.name}（参考价 ¥{_num(product.reference_price):.2f}）共 {len(rows)} 条报价，全国农产品近期均价 ¥{avg_price if avg_price else '?'}。",
    }


async def _xinfadi_rows(db: AsyncSession, product_name: str = "", limit: int = 500) -> list[dict[str, Any]]:
    product_name = _extract_xinfadi_keyword(product_name)
    clause = "WHERE product_name LIKE :p" if product_name else ""
    params = {"limit": limit}
    if product_name:
        params["p"] = f"%{product_name}%"
    result = await db.execute(
        text(
            f"""
            SELECT crawl_date, product_name, category1, min_price, avg_price, max_price, unit, origin
            FROM `{settings.xinfadi_price_table}`
            {clause}
            ORDER BY crawl_date DESC
            LIMIT :limit
            """
        ),
        params,
    )
    return [
        {
            "date": row.crawl_date.isoformat() if row.crawl_date else "",
            "product_name": _repair_text(row.product_name),
            "category": _repair_text(row.category1),
            "min_price": _num(row.min_price),
            "avg_price": _num(row.avg_price),
            "max_price": _num(row.max_price),
            "unit": _repair_text(row.unit),
            "origin": _repair_text(row.origin),
        }
        for row in result.fetchall()
    ]


async def _xinfadi_product_names(db: AsyncSession, query: str, limit: int = 20) -> list[str]:
    keyword = (query or "").strip()
    if not keyword:
        return []
    result = await db.execute(
        text(
            f"""
            SELECT product_name, COUNT(*) AS sample_count
            FROM `{settings.xinfadi_price_table}`
            WHERE product_name LIKE :p
            GROUP BY product_name
            ORDER BY sample_count DESC, product_name ASC
            LIMIT :limit
            """
        ),
        {"p": f"%{keyword}%", "limit": int(limit)},
    )
    return [_repair_text(row.product_name) for row in result.fetchall()]


def _load_json(value: Any, fallback: Any) -> Any:
    if value is None or value == "":
        return fallback
    if isinstance(value, (dict, list)):
        return value
    try:
        return json.loads(str(value))
    except Exception:
        return fallback


def _normalize_forecast_row(row: dict[str, Any]) -> dict[str, Any]:
    yhat = _num(row.get("yhat"))
    lower = _num(row.get("yhat_lower", row.get("lower")))
    upper = _num(row.get("yhat_upper", row.get("upper")))
    return {
        "date": str(row.get("date") or ""),
        "yhat": round(yhat, 4),
        "price": round(yhat, 4),
        "yhat_lower": round(lower, 4),
        "yhat_upper": round(upper, 4),
        "lower": round(lower, 4),
        "upper": round(upper, 4),
        "confidence": round(_num(row.get("confidence")), 4),
        "trend": row.get("trend") or "flat",
    }


def _forecast_rows_from_tomorrow(ensemble: list[dict[str, Any]], days: int, start_day: Optional[date] = None) -> list[dict[str, Any]]:
    tomorrow = start_day or (_today() + timedelta(days=1))
    future = []
    for row in ensemble:
        try:
            row_day = date.fromisoformat(str(row.get("date") or "")[:10])
        except ValueError:
            continue
        if row_day >= tomorrow:
            future.append(row)
    return future[: max(1, min(30, int(days or 7)))]


def _forecast(rows: list[dict[str, Any]], days: int) -> list[dict[str, Any]]:
    values = [_num(row["avg_price"]) for row in rows if _num(row.get("avg_price")) > 0]
    baseline = sum(values[:14]) / len(values[:14]) if values[:14] else 1
    vol = 0.04
    if len(values) > 2 and baseline:
        vol = min(0.22, max(0.015, (max(values[:30]) - min(values[:30])) / baseline))
    last = _today()
    if rows and rows[0].get("date"):
        try:
            last = date.fromisoformat(rows[0]["date"])
        except ValueError:
            pass
    out = []
    for i in range(1, days + 1):
        drift = ((i % 7) - 3) * vol * 0.15
        yhat = baseline * (1 + drift)
        out.append({"date": (last + timedelta(days=i)).isoformat(), "price": round(yhat, 2), "lower": round(yhat * (1 - vol), 2), "upper": round(yhat * (1 + vol), 2)})
    return out


async def _tool_national_ag_price(db: AsyncSession, args: dict[str, Any]) -> dict[str, Any]:
    raw = str(args.get("query_text") or args.get("product_name") or "").strip()
    pq_src = str(args.get("product_query") or raw).strip()
    parsed = national_price_svc.parse_price_query(pq_src if args.get("product_query") else raw)
    if args.get("product_query"):
        parsed = national_price_svc.ParsedPriceQuery(
            raw_text=raw or pq_src,
            product_query=str(args.get("product_query") or parsed.product_query),
            goods_name_hint=parsed.goods_name_hint or national_price_svc.parse_price_query(str(args["product_query"])).goods_name_hint,
            grade_hint=parsed.grade_hint,
            target_date=parsed.target_date,
            intent=parsed.intent,
        )
    td: Optional[date] = parsed.target_date
    if args.get("target_date"):
        try:
            td = date.fromisoformat(str(args["target_date"])[:10])
        except ValueError:
            td = parsed.target_date
    return await national_price_svc.history_for_query(db, raw, parsed=parsed, target_date=td)


async def _tool_national_ag_forecast(db: AsyncSession, args: dict[str, Any]) -> dict[str, Any]:
    mode = str(args.get("mode") or "future").strip().lower()
    if mode not in {"tomorrow", "future"}:
        mode = "future"
    target_offset = max(1, min(30, int(args.get("target_offset") or 1)))
    raw_text = str(args.get("query_text") or args.get("product_name") or "").strip()
    parsed = national_price_svc.parse_price_query(raw_text)
    days = max(1, min(30, int(args.get("days") or 7)))
    return await national_price_svc.forecast_for_query(
        db,
        query_text=raw_text,
        mode=mode,
        days=days,
        target_offset=target_offset,
        parsed=parsed,
    )


async def _run_national_price_tool(
    db: AsyncSession,
    name: str,
    args: dict[str, Any],
    *,
    on_training: Optional[Any] = None,
) -> dict[str, Any]:
    """执行全国价格工具；预测类在缺快照时自动训练。"""
    if name in ("get_national_ag_forecast_price", "get_xinfadi_forecast_price"):
        raw_text = str(args.get("query_text") or "").strip()
        parsed = national_price_svc.parse_price_query(raw_text)
        sku_key, _, _ = await national_price_svc.resolve_sku_from_parsed(db, parsed)
        if sku_key and not await zg_train_svc.snapshot_exists(db, sku_key):
            async for ev in zg_train_svc.iter_forecast_train_progress(db, sku_key):
                if on_training:
                    on_training(ev)
        return await _tool_national_ag_forecast(db, args)
    return await _tool_national_ag_price(db, args)


async def _tool_xinfadi_price(db: AsyncSession, args: dict[str, Any]) -> dict[str, Any]:
    return await _tool_national_ag_price(db, args)


async def _tool_xinfadi_forecast(db: AsyncSession, args: dict[str, Any]) -> dict[str, Any]:
    return await _tool_national_ag_forecast(db, args)




async def _tool_lookup_entity_by_name(db: AsyncSession, args: dict[str, Any]) -> dict[str, Any]:
    """模糊匹配实体名（中文公司名 / 用户名 / 食堂名），把名字映射成 ID + 类型。
    用户问『中农食迅的日报』时，AI 应先调用此工具把"中农食迅"映射成 delivery_id=3。"""
    name = (args.get("name") or "").strip()
    if not name:
        return {"type": "error", "rows": [], "summary": "缺少 name 参数。"}
    role = (args.get("role") or "").strip().lower() or None
    limit = _clamp_limit(args, default=5, cap=10)

    out: list[dict[str, Any]] = []
    # 1. 模糊匹配 users（按 company_name 或 username）
    user_stmt = select(User).where(
        (User.company_name.like(f"%{name}%")) | (User.username.like(f"%{name}%"))
    )
    if role in {"client", "delivery", "supplier", "factory"}:
        user_stmt = user_stmt.where(User.role == role)
    users = (await db.scalars(user_stmt.limit(limit))).all()
    for u in users:
        out.append({
            "type": u.role,
            "id": u.id,
            "name": u.company_name or u.username,
            "username": u.username,
            "role": u.role,
            "match_field": "company_name" if u.company_name and name in u.company_name else "username",
        })
    # 2. 模糊匹配 client_canteens（按 name）
    if role in (None, "canteen"):
        canteens = (
            await db.scalars(select(ClientCanteen).where(ClientCanteen.name.like(f"%{name}%")).limit(limit))
        ).all()
        if canteens:
            school_ids = {int(c.school_client_id) for c in canteens}
            school_map = await _user_name_map(db, school_ids)
            for c in canteens:
                out.append({
                    "type": "canteen",
                    "id": c.id,
                    "name": c.name,
                    "school_client_id": c.school_client_id,
                    "school_name": school_map.get(int(c.school_client_id), ""),
                    "match_field": "canteen_name",
                })
    return {
        "type": "entity_lookup",
        "title": f"实体匹配：{name}",
        "rows": out[:limit],
        "summary": f"匹配到 {len(out)} 个候选" + (
            f"，其中最佳：{out[0]['type']} «{out[0]['name']}» (id={out[0]['id']})" if out else "。请确认名称"
        ),
    }


async def _tool_schema_overview(db: AsyncSession, _args: dict[str, Any]) -> dict[str, Any]:
    result = await db.execute(text("SHOW TABLES"))
    tables = [list(row._mapping.values())[0] for row in result.fetchall()]
    return {
        "type": "plain",
        "title": "可查询数据范围",
        "rows": [{"table": t} for t in tables if any(k in t for k in ("order", "alert", "xinfadi", "product", "delivery", "bill"))][:50],
        "summary": (
            "我能回答两类问题：\n"
            "① 本系统监管数据 —— 订单 / 账单对账·超期 / 工单投诉 / 质检 / 履约完成率 / 配送商·供货商·学校指标 / 设备 IoT / 告警 / 财务报表 / 日周月报。"
            "例：『中农食迅这个月履约率』『超期未付账单』『给我生成财务报表』。\n"
            "② 全国农产品价格 —— 历史菜价、明天/未来价格预测（与数据挖掘中心全国 Tab 一致）。例：『大白菜价格』『大白菜明天多少钱』。\n"
            "换成上面这样的问法，我就能给你真实数据。"
        ),
    }


def _report_period_from_type(report_type: str) -> tuple[date, date, str]:
    rt = (report_type or "daily").strip().lower()
    today = _today()
    if rt in ("financial", "财务", "财务报表", "经营"):
        return today.replace(day=1), today, "财务报表"
    if rt in ("weekly", "week", "周报"):
        return today - timedelta(days=6), today, "周报"
    if rt in ("monthly", "month", "月报"):
        return today.replace(day=1), today, "月报"
    return today, today, "日报"


async def _statement_finance_summary(db: AsyncSession, start: date, end: date) -> dict[str, Any]:
    """对 BillingStatement 在区间内做财务聚合（按方向 SUM），不复用 50 行上限的 search。"""
    from sqlalchemy import func as _func

    start_dt = datetime.combine(start, datetime.min.time())
    end_dt = datetime.combine(end + timedelta(days=1), datetime.min.time())
    grouped = (
        await db.execute(
            select(
                BillingStatement.direction,
                _func.count().label("cnt"),
                _func.coalesce(_func.sum(BillingStatement.amount), 0).label("amount"),
                _func.coalesce(_func.sum(BillingStatement.settled_amount), 0).label("settled"),
            )
            .where(BillingStatement.created_at >= start_dt, BillingStatement.created_at < end_dt)
            .group_by(BillingStatement.direction)
        )
    ).all()
    agg: dict[str, dict[str, Any]] = {
        "应付": {"cnt": 0, "amount": 0.0, "settled": 0.0, "unsettled": 0.0},
        "应收": {"cnt": 0, "amount": 0.0, "settled": 0.0, "unsettled": 0.0},
    }
    for r in grouped:
        d = str(r.direction or "")
        bucket = agg.setdefault(d, {"cnt": 0, "amount": 0.0, "settled": 0.0, "unsettled": 0.0})
        bucket["cnt"] = int(r.cnt or 0)
        bucket["amount"] = _num(r.amount)
        bucket["settled"] = _num(r.settled)
    for v in agg.values():
        v["unsettled"] = round(max(0.0, _num(v["amount"]) - _num(v["settled"])), 2)
    top_rows = (
        await db.scalars(
            select(BillingStatement)
            .where(BillingStatement.created_at >= start_dt, BillingStatement.created_at < end_dt)
            .order_by(BillingStatement.amount.desc())
            .limit(10)
        )
    ).all()
    uid = {int(r.owner_user_id) for r in top_rows} | {int(r.counterparty_user_id) for r in top_rows}
    umap = await _user_name_map(db, uid)
    top = [
        {
            "statement_no": r.statement_no,
            "direction": r.direction,
            "status": r.status,
            "amount": _num(r.amount),
            "counterparty_name": umap.get(int(r.counterparty_user_id), ""),
        }
        for r in top_rows
    ]
    return {"by_direction": agg, "top": top}


async def _resolve_subject(db: AsyncSession, args: dict[str, Any]) -> tuple[Optional[str], Optional[int], Optional[str], list[dict[str, Any]]]:
    """解析报告主体：返回 (subject_type, subject_id, subject_name, candidates)。
    优先用 subject_id+subject_type；否则用 subject_name 走模糊匹配；都没给返回 (None, None, None, [])。"""
    sid = args.get("subject_id")
    stype = (args.get("subject_type") or "").strip().lower() or None
    sname_in = (args.get("subject_name") or "").strip()
    if sid and stype:
        if stype == "canteen":
            c = await db.get(ClientCanteen, int(sid))
            if c:
                return "canteen", int(c.id), c.name, []
        else:
            u = await db.get(User, int(sid))
            if u:
                return stype, int(u.id), u.company_name or u.username, []
    if sname_in:
        lookup = await _tool_lookup_entity_by_name(db, {"name": sname_in, "limit": 5})
        rows = lookup.get("rows") or []
        if rows:
            top = rows[0]
            return top["type"], int(top["id"]), top["name"], rows
        return None, None, None, []
    return None, None, None, []


def _md_table(headers: list[str], rows: list[list[Any]]) -> str:
    if not rows:
        return "_（无记录）_\n"
    out = "| " + " | ".join(headers) + " |\n"
    out += "| " + " | ".join(["---"] * len(headers)) + " |\n"
    for r in rows:
        out += "| " + " | ".join(str(x if x is not None else "") for x in r) + " |\n"
    return out


def _order_status_rows(orders: list[Any]) -> list[list[Any]]:
    bucket: dict[str, int] = {}
    for order in orders:
        status = str(getattr(order, "status", "") or "未知")
        bucket[status] = bucket.get(status, 0) + 1
    return [[status, count] for status, count in sorted(bucket.items(), key=lambda item: -item[1])]


def _alert_type_rows(alerts: list[dict[str, Any]]) -> list[list[Any]]:
    bucket: dict[str, int] = {}
    for alert in alerts:
        alert_type = str(alert.get("type") or "unknown")
        bucket[alert_type] = bucket.get(alert_type, 0) + 1
    return [[alert_type, count] for alert_type, count in sorted(bucket.items(), key=lambda item: -item[1])]


def _global_report_suggestions(
    kpi_rows: list[dict[str, Any]],
    alert_type_rows: list[list[Any]],
    overdue_rows: list[dict[str, Any]],
    ticket_rows: list[dict[str, Any]],
    orders: list[Any],
) -> list[str]:
    suggestions: list[str] = []
    if alert_type_rows:
        top_type, top_count = alert_type_rows[0]
        suggestions.append(f"优先关注当日最多的告警类型「{top_type}」（{top_count} 条）。")
    if overdue_rows:
        suggestions.append(f"仍有 {len(overdue_rows)} 张超期未付账单，建议催办对账。")
    fulfilled_rate = next((float(r.get("value") or 0) for r in kpi_rows if r.get("metric") == "履约完成率"), 0.0)
    if orders and fulfilled_rate < 80:
        suggestions.append(f"履约完成率 {fulfilled_rate}%，低于 80%，需跟踪未确认收货订单。")
    pending_tickets = sum(1 for ticket in ticket_rows if ticket.get("status") == "待处理")
    if ticket_rows:
        suggestions.append(f"当日新增工单 {len(ticket_rows)} 条，其中 {pending_tickets} 条待处理。")
    if not suggestions:
        suggestions.append("当日业务整体平稳，建议继续例行巡查。")
    return suggestions


async def _tool_generate_report(db: AsyncSession, args: dict[str, Any]) -> dict[str, Any]:
    """生成日报/周报/月报。支持主体过滤：传 subject_name（自然语言名字）或 subject_id+subject_type。"""
    start, end, period_label = _report_period_from_type(str(args.get("report_type") or "daily"))
    if args.get("start_date"):
        s = _parse_date(args.get("start_date"))
        if s:
            start = s
    if args.get("end_date"):
        e = _parse_date(args.get("end_date"))
        if e:
            end = e
    if end < start:
        start, end = end, start

    report_type_raw = str(args.get("report_type") or "daily").strip().lower()
    subject_name_raw = (args.get("subject_name") or "").strip()
    is_financial = report_type_raw in ("financial", "财务", "财务报表", "经营") or any(
        w in subject_name_raw for w in _FINANCIAL_WORDS
    )
    # 概念词（财务/经营/统计/汇总…）不是实体名：剥掉后若为空，降级为全局报告，不再硬失败
    if subject_name_raw:
        stripped = subject_name_raw
        for w in _NON_ENTITY_WORDS:
            stripped = stripped.replace(w, "")
        if not stripped.strip():
            args = {**args, "subject_name": ""}

    stype, sid, sname, candidates = await _resolve_subject(db, args)

    # ===== 全局财务报表（未指定主体）=====
    if not stype and is_financial:
        kpi = await _tool_kpi_summary(db, {"scope": "range", "start_date": start.isoformat(), "end_date": end.isoformat()})
        fin = await _statement_finance_summary(db, start, end)
        overdue = await _tool_get_overdue_bills(db, {"scope": "payment", "limit": 10})
        pay = fin["by_direction"].get("应付", {})
        rec = fin["by_direction"].get("应收", {})
        report = (
            f"# 财务报表（{start} ~ {end}）\n\n"
            f"## 经营概览\n\n{kpi['summary']}\n\n"
            "## 应付 / 应收\n\n"
            + _md_table(
                ["科目", "笔数", "金额", "已结", "未结"],
                [
                    ["应付", pay.get("cnt", 0), f"¥{_num(pay.get('amount')):.2f}", f"¥{_num(pay.get('settled')):.2f}", f"¥{_num(pay.get('unsettled')):.2f}"],
                    ["应收", rec.get("cnt", 0), f"¥{_num(rec.get('amount')):.2f}", f"¥{_num(rec.get('settled')):.2f}", f"¥{_num(rec.get('unsettled')):.2f}"],
                ],
            )
            + f"\n## 超期未付账单（{len(overdue['rows'])} 张）\n\n"
            + _md_table(["账单号", "对方", "金额", "超期天数"], [[r["statement_no"], r.get("counterparty_name", ""), f"¥{_num(r['amount']):.2f}", r.get("overdue_days", 0)] for r in overdue["rows"][:10]])
            + "\n## 金额 Top 账单\n\n"
            + _md_table(["账单号", "方向", "状态", "金额", "对方"], [[r["statement_no"], r["direction"], r["status"], f"¥{_num(r['amount']):.2f}", r["counterparty_name"]] for r in fin["top"][:10]])
            + "\n## 监管建议\n\n- 关注应收未结回款进度，及时催办对账。\n- 优先清理超期未付账单，规避履约违约风险。\n"
        )
        summary = (
            f"已生成 {start}~{end} 财务报表：应付 ¥{_num(pay.get('amount')):.2f}（未结 ¥{_num(pay.get('unsettled')):.2f}），"
            f"应收 ¥{_num(rec.get('amount')):.2f}（未结 ¥{_num(rec.get('unsettled')):.2f}），超期未付 {len(overdue['rows'])} 张。"
        )
        return {"type": "report", "title": f"财务报表（{start}~{end}）", "rows": [], "summary": summary, "report_content": report}

    # 用户给了 subject_name 但匹配不到 → 直接告诉用户，不降级到全局日报
    if not stype and (args.get("subject_name") or "").strip():
        raw_name = args.get("subject_name")
        # 尝试一次更宽松的搜索给出建议
        lookup_loose = await _tool_lookup_entity_by_name(db, {"name": raw_name[0] if raw_name else "", "limit": 5})
        suggestion = "、".join(f"{r['name']}({r['type']})" for r in (lookup_loose.get("rows") or [])[:5])
        return {
            "type": "notice",
            "title": "未找到主体",
            "rows": [],
            "summary": (
                f"未匹配到名为 «{raw_name}» 的客户/食堂/配送商/供货商/厂家。"
                + (f" 可能你想问：{suggestion}？" if suggestion else " 请确认名称或先用 lookup_entity_by_name 查找候选。")
            ),
        }
    ambiguous_hint = ""

    # ===== 全局日报（未指定主体）=====
    if not stype:
        range_args = {"scope": "range", "start_date": start.isoformat(), "end_date": end.isoformat()}
        kpi = await _tool_kpi_summary(db, range_args)
        orders = await _orders_in_range(db, start, end)
        region_rank = await _tool_rank(db, {**range_args, "limit": 10}, "region")
        category_rank = await _tool_rank(db, {**range_args, "limit": 10}, "category")
        alerts = await _tool_search_alerts(
            db,
            {"start_date": start.isoformat(), "end_date": end.isoformat(), "limit": 30},
        )
        overdue = await _tool_get_overdue_bills(db, {"scope": "payment", "limit": 20})
        tickets = await _tool_search_tickets(
            db,
            {"start_date": start.isoformat(), "end_date": end.isoformat(), "limit": 15},
        )
        alert_types = _alert_type_rows(alerts["rows"])
        suggestions = _global_report_suggestions(kpi["rows"], alert_types, overdue["rows"], tickets["rows"], orders)
        report = (
            f"# 监管{period_label}（{start} ~ {end}）\n\n"
            + (f"> {ambiguous_hint}\n\n" if ambiguous_hint else "")
            + "## 核心 KPI\n\n"
            + _md_table(["指标", "数值", "单位"], [[r["metric"], r["value"], r["unit"]] for r in kpi["rows"]])
            + "\n## 订单状态分布\n\n"
            + _md_table(["状态", "订单数"], _order_status_rows(orders))
            + "\n## 区域排行 Top10\n\n"
            + _md_table(
                ["区域", "订单数", "GMV"],
                [[r["name"], r["order_count"], r["gmv"]] for r in region_rank["rows"]],
            )
            + "\n## 品类分布 Top10\n\n"
            + _md_table(
                ["品类", "订单数", "金额"],
                [[r["name"], r["order_count"], r["amount"]] for r in category_rank["rows"]],
            )
            + f"\n## 当日告警类型汇总（{len(alerts['rows'])} 条）\n\n"
            + _md_table(["类型", "数量"], alert_types)
            + f"\n## 当日告警明细（{len(alerts['rows'])} 条）\n\n"
            + _md_table(
                ["等级", "类型", "描述"],
                [[a["level"], a["type"], a["description"]] for a in alerts["rows"][:30]],
            )
            + f"\n## 超期未付账单（{len(overdue['rows'])} 张）\n\n"
            + _md_table(
                ["账单号", "对方", "金额", "超期天数"],
                [
                    [r["statement_no"], r.get("counterparty_name", ""), r["amount"], r.get("overdue_days", 0)]
                    for r in overdue["rows"][:20]
                ],
            )
            + f"\n## 当日工单（{len(tickets['rows'])} 条）\n\n"
            + _md_table(
                ["工单ID", "订单号", "类型", "状态"],
                [[t["id"], t.get("order_no", ""), t["type"], t["status"]] for t in tickets["rows"][:15]],
            )
            + "\n## 监管建议\n\n"
            + "\n".join(f"- {item}" for item in suggestions)
            + "\n"
        )
        summary = (
            f"已生成 {start}~{end} 监管{period_label}。"
            f"{kpi['summary']}"
        )
        return {"type": "report", "title": f"监管{period_label}", "rows": [], "summary": summary, "report_content": report}

    # ===== 主体报告 =====
    parts: list[str] = [f"# {sname} {period_label}（{start} ~ {end}）\n"]
    parts.append(f"_主体类型：{stype}，ID={sid}_\n\n")

    if stype == "delivery":
        metrics = await _tool_get_delivery_metrics(db, {"delivery_id": sid, "start_date": start.isoformat(), "end_date": end.isoformat()})
        orders = await _tool_search_orders(db, {"delivery_id": sid, "start_date": start.isoformat(), "end_date": end.isoformat(), "limit": 30})
        tickets = await _tool_search_tickets(db, {"start_date": start.isoformat(), "end_date": end.isoformat(), "limit": 10})
        rel_tickets = [t for t in tickets["rows"] if t.get("assigned_delivery_id") == sid or any(o.get("order_id") == t.get("order_id") for o in orders["rows"])]
        parts.append("## 履约指标\n\n" + _md_table(["指标", "数值", "单位"], [[r["metric"], r["value"], r["unit"]] for r in metrics["rows"]]))
        parts.append("\n## 客户分布 Top5\n\n" + _md_table(["客户", "订单数", "GMV"], [[r["client_name"], r["order_count"], r["gmv"]] for r in metrics.get("top_clients", [])]))
        parts.append(f"\n## 订单清单（{len(orders['rows'])} 单）\n\n" + _md_table(["订单号", "客户", "食堂", "状态", "金额"], [[o["order_no"], o["client_name"], o["canteen_name"], o["status"], o["amount"]] for o in orders["rows"][:30]]))
        parts.append(f"\n## 相关工单（{len(rel_tickets)} 条）\n\n" + _md_table(["工单ID", "订单号", "类型", "状态"], [[t["id"], t.get("order_no", ""), t["type"], t["status"]] for t in rel_tickets[:10]]))
    elif stype == "client":
        metrics = await _tool_get_client_metrics(db, {"client_id": sid, "start_date": start.isoformat(), "end_date": end.isoformat()})
        orders = await _tool_search_orders(db, {"client_id": sid, "start_date": start.isoformat(), "end_date": end.isoformat(), "limit": 30})
        parts.append("## 采购指标\n\n" + _md_table(["指标", "数值", "单位"], [[r["metric"], r["value"], r["unit"]] for r in metrics["rows"]]))
        parts.append("\n## 按食堂拆分\n\n" + _md_table(["食堂", "订单数", "GMV"], [[r["canteen_name"], r["order_count"], r["gmv"]] for r in metrics.get("by_canteen", [])]))
        parts.append("\n## 品类 Top10\n\n" + _md_table(["品类", "金额", "数量"], [[r["category_name"], r["amount"], r["qty"]] for r in metrics.get("by_category", [])]))
        parts.append("\n## 配送商分布\n\n" + _md_table(["配送商", "订单数", "GMV"], [[r["delivery_name"], r["order_count"], r["gmv"]] for r in metrics.get("by_delivery", [])]))
        parts.append(f"\n## 订单清单（{len(orders['rows'])} 单）\n\n" + _md_table(["订单号", "食堂", "配送商", "状态", "金额"], [[o["order_no"], o["canteen_name"], o["delivery_name"], o["status"], o["amount"]] for o in orders["rows"][:30]]))
    elif stype == "canteen":
        # 食堂日报：通过 orders.canteen_id 过滤
        orders = await _tool_search_orders(db, {"canteen_id": sid, "start_date": start.isoformat(), "end_date": end.isoformat(), "limit": 50})
        stmts = await _tool_search_statements(db, {"canteen_id": sid, "start_date": start.isoformat(), "end_date": end.isoformat(), "limit": 20})
        gmv = sum(_num(o["amount"]) for o in orders["rows"])
        parts.append(f"## 食堂采购汇总\n\n- 订单数：**{len(orders['rows'])}** 单\n- GMV：**¥{gmv:.2f}**\n- 关联账单：{len(stmts['rows'])} 张\n\n")
        parts.append(f"## 订单清单（{len(orders['rows'])} 单）\n\n" + _md_table(["订单号", "配送商", "状态", "金额"], [[o["order_no"], o["delivery_name"], o["status"], o["amount"]] for o in orders["rows"][:30]]))
        parts.append(f"\n## 应付账单（{len(stmts['rows'])} 张）\n\n" + _md_table(["账单号", "方向", "状态", "金额", "已结"], [[r["statement_no"], r["direction"], r["status"], r["amount"], r["settled_amount"]] for r in stmts["rows"]]))
    elif stype == "supplier":
        metrics = await _tool_get_supplier_metrics(db, {"supplier_id": sid, "start_date": start.isoformat(), "end_date": end.isoformat()})
        quotes = await _tool_get_supplier_quotes(db, {"supplier_id": sid, "limit": 20})
        qr = await _tool_search_quality_reports(db, {"supplier_id": sid, "start_date": start.isoformat(), "end_date": end.isoformat(), "limit": 10})
        parts.append("## 供货指标\n\n" + _md_table(["指标", "数值", "单位"], [[r["metric"], r["value"], r["unit"]] for r in metrics["rows"]]))
        parts.append("\n## 商品 Top10\n\n" + _md_table(["商品", "数量", "金额"], [[r["product_name"], r["qty"], r["amount"]] for r in metrics.get("top_products", [])]))
        parts.append(f"\n## 报价清单（{len(quotes['rows'])} 条）\n\n" + _md_table(["商品", "报价"], [[r["product_name"], r["quote_price"]] for r in quotes["rows"][:20]]))
        parts.append(f"\n## 质检报告（{len(qr['rows'])} 条）\n\n" + _md_table(["报告号", "商品", "状态"], [[r["report_no"], r.get("product_name", ""), r["status"]] for r in qr["rows"]]))
    elif stype == "factory":
        metrics = await _tool_get_factory_metrics(db, {"factory_id": sid, "start_date": start.isoformat(), "end_date": end.isoformat()})
        parts.append("## 厂家指标\n\n" + _md_table(["指标", "数值", "单位"], [[r["metric"], r["value"], r["unit"]] for r in metrics["rows"]]))
    else:
        parts.append(f"_未支持的主体类型：{stype}_\n")

    report = "".join(parts)
    summary = f"已生成 {sname}（{stype}）的{period_label}：{start}~{end}。"
    if len(candidates) > 1:
        seen_alt: set[str] = set()
        alt_parts: list[str] = []
        for c in candidates[1:]:
            label = c["name"]
            if c.get("type") == "canteen" and c.get("school_name"):
                label = f"{c['school_name']}-{c['name']}"
            key = f"{c['type']}:{c['id']}"
            if key in seen_alt:
                continue
            seen_alt.add(key)
            alt_parts.append(f"{label}({c['type']})")
        if alt_parts:
            summary += f" 备选候选：{'、'.join(alt_parts[:3])}。"
    return {
        "type": "report",
        "title": f"{sname} {period_label}",
        "rows": [],
        "subject": {"type": stype, "id": sid, "name": sname},
        "candidates": candidates,
        "summary": summary,
        "report_content": report,
    }


async def _tool_search_docs(db: AsyncSession, args: dict[str, Any]) -> dict[str, Any]:
    _ = db
    query = str(args.get("query") or args.get("q") or "").strip()
    limit = max(1, min(10, int(args.get("limit") or 5)))
    hits = search_docs(query, top_k=limit, audience="monitor")
    citations = citations_from_rag(hits)
    rows = [
        {
            "doc": h.get("doc_title"),
            "section": h.get("section"),
            "excerpt": ((h.get("text") or "")[:200] + "…") if len(h.get("text") or "") > 200 else (h.get("text") or ""),
            "score": h.get("score"),
        }
        for h in hits
    ]
    return {
        "type": "docs",
        "chart_type": "none",
        "title": "操作手册检索",
        "columns": [
            {"key": "doc", "label": "文档"},
            {"key": "section", "label": "章节"},
            {"key": "excerpt", "label": "摘录"},
        ],
        "rows": rows,
        "citations": citations,
        "summary": f"检索到 {len(hits)} 条相关手册摘录。" if hits else "未找到匹配的操作手册段落，请换关键词。",
    }


TOOLS: list[dict[str, Any]] = [
    # ===== 基础 KPI / 排行 / 趋势 =====
    {"type": "function", "function": {"name": "get_kpi_summary", "description": "查询核心 KPI（订单数/GMV/当日告警/履约完成率）。when_to_use: 问『生意咋样/概况/今天多少单』等概览。when_not_to_use: 不要用于订单明细列表（用 search_orders）或复杂聚合排行（用 run_sql）。日期：昨天/今天传 start_date+end_date ISO。", "parameters": {"type": "object", "properties": {"scope": {"type": "string", "enum": ["today", "range"]}, "start_date": {"type": "string"}, "end_date": {"type": "string"}, "client_id": {"type": "integer"}, "canteen_id": {"type": "integer"}, "delivery_id": {"type": "integer"}, "supplier_id": {"type": "integer"}}, "required": ["scope"]}}},
    {"type": "function", "function": {"name": "get_region_rank", "description": "查询区域 GMV 排行。", "parameters": {"type": "object", "properties": {"start_date": {"type": "string"}, "end_date": {"type": "string"}, "limit": {"type": "integer"}}}}},
    {"type": "function", "function": {"name": "get_category_distribution", "description": "查询一级品类销售分布（订单数+金额）。用户要饼图/占比时 chart_type 传 pie。", "parameters": {"type": "object", "properties": {"start_date": {"type": "string"}, "end_date": {"type": "string"}, "limit": {"type": "integer"}, "chart_type": {"type": "string", "enum": ["bar", "pie"], "description": "饼图用 pie，柱状排行用 bar"}}, "required": []}}},
    {"type": "function", "function": {"name": "get_top_goods", "description": "查询商品销售排行。", "parameters": {"type": "object", "properties": {"start_date": {"type": "string"}, "end_date": {"type": "string"}, "limit": {"type": "integer"}}}}},
    {"type": "function", "function": {"name": "get_daily_trend", "description": "查询每日 GMV 和订单趋势。", "parameters": {"type": "object", "properties": {"start_date": {"type": "string"}, "end_date": {"type": "string"}}}}},
    {"type": "function", "function": {"name": "get_today_orders", "description": "查询今日订单清单（含客户/食堂/配送商名）。", "parameters": {"type": "object", "properties": {}}}},
    # ===== 类 A：实体溯源（按 ID/编号查详情）=====
    {"type": "function", "function": {"name": "get_order_detail", "description": "按订单号或订单 ID 查询完整订单详情：客户、食堂、配送商、供货商、分单明细、关联账单、工单、异常。用户问『这个单/该订单是什么/谁下的/详情』时优先调用。", "parameters": {"type": "object", "properties": {"order_no": {"type": "string", "description": "订单号，如 OD180253..."}, "order_id": {"type": "integer"}}}}},
    {"type": "function", "function": {"name": "get_statement_detail", "description": "按账单号或 ID 查询账单详情，含 owner/counterparty/食堂/关联订单。", "parameters": {"type": "object", "properties": {"statement_no": {"type": "string"}, "statement_id": {"type": "integer"}}}}},
    {"type": "function", "function": {"name": "get_ticket_detail", "description": "按工单 ID 查询工单详情：类型、状态、处理流水、关联订单。", "parameters": {"type": "object", "properties": {"ticket_id": {"type": "integer"}}, "required": ["ticket_id"]}}},
    {"type": "function", "function": {"name": "get_contract_detail", "description": "按合约号或 ID 查询合约详情：买卖双方、有效期、品类、相关食堂。", "parameters": {"type": "object", "properties": {"contract_no": {"type": "string"}, "contract_id": {"type": "integer"}}}}},
    {"type": "function", "function": {"name": "get_user_detail", "description": "按用户名或用户 ID 查询用户详情。client 返回所属食堂；delivery 返回车辆/设备数/绑定供货商；supplier 返回报价数。", "parameters": {"type": "object", "properties": {"username": {"type": "string"}, "user_id": {"type": "integer"}}}}},
    # ===== 类 B：列表搜索 =====
    {"type": "function", "function": {"name": "search_orders", "description": "按多条件搜索订单（客户/食堂/配送商/供货商/状态/异常/金额/时间）。返回含客户名/食堂名/配送商名。默认最近 7 天。when_not_to_use: 只要 KPI 概览数字时用 get_kpi_summary，不要拉全量明细。", "parameters": {"type": "object", "properties": {"client_id": {"type": "integer"}, "canteen_id": {"type": "integer"}, "delivery_id": {"type": "integer"}, "supplier_id": {"type": "integer"}, "status": {"type": "string", "description": "下单/配货/发货/收货/收货确认/已结算/取消"}, "has_abnormal": {"type": "boolean"}, "amount_min": {"type": "number"}, "amount_max": {"type": "number"}, "order_no_like": {"type": "string"}, "start_date": {"type": "string"}, "end_date": {"type": "string"}, "limit": {"type": "integer"}}}}},
    {"type": "function", "function": {"name": "search_statements", "description": "按多条件搜索账单（应付/应收、状态、所有者、食堂、时间）。", "parameters": {"type": "object", "properties": {"direction": {"type": "string", "enum": ["应付", "应收"]}, "status": {"type": "string", "description": "待确认/已确认/部分结清/已结清"}, "owner_user_id": {"type": "integer"}, "counterparty_user_id": {"type": "integer"}, "canteen_id": {"type": "integer"}, "role": {"type": "string", "enum": ["client", "delivery", "supplier", "factory"]}, "start_date": {"type": "string"}, "end_date": {"type": "string"}, "limit": {"type": "integer"}}}}},
    {"type": "function", "function": {"name": "search_tickets", "description": "按多条件搜索工单：状态、类型、订单 ID、时间。", "parameters": {"type": "object", "properties": {"status": {"type": "string", "enum": ["待处理", "处理中", "已关闭"]}, "type": {"type": "string", "enum": ["异常订单", "售后投诉", "配送异常"]}, "order_id": {"type": "integer"}, "start_date": {"type": "string"}, "end_date": {"type": "string"}, "limit": {"type": "integer"}}}}},
    {"type": "function", "function": {"name": "search_quality_reports", "description": "按多条件搜索质检报告。", "parameters": {"type": "object", "properties": {"supplier_id": {"type": "integer"}, "order_id": {"type": "integer"}, "product_id": {"type": "integer"}, "status": {"type": "string", "enum": ["待审核", "已通过"]}, "start_date": {"type": "string"}, "end_date": {"type": "string"}, "limit": {"type": "integer"}}}}},
    {"type": "function", "function": {"name": "search_alerts", "description": "按多条件搜索告警（等级/状态/类型/时间）。", "parameters": {"type": "object", "properties": {"level": {"type": "string", "enum": ["high", "medium", "low"]}, "status": {"type": "string", "enum": ["open", "closed"]}, "type": {"type": "string"}, "start_date": {"type": "string"}, "end_date": {"type": "string"}, "limit": {"type": "integer"}}}}},
    # 兼容旧名：get_ops_alerts 仍指向 search_alerts
    {"type": "function", "function": {"name": "get_ops_alerts", "description": "（兼容）查询运营预警。建议改用 search_alerts。", "parameters": {"type": "object", "properties": {"limit": {"type": "integer"}}}}},
    # ===== 类 C：主体聚合 =====
    {"type": "function", "function": {"name": "get_delivery_metrics", "description": "查询配送商履约指标：单数/GMV/完成率/按期率/异常率/客户分布 top5。", "parameters": {"type": "object", "properties": {"delivery_id": {"type": "integer"}, "start_date": {"type": "string"}, "end_date": {"type": "string"}}, "required": ["delivery_id"]}}},
    {"type": "function", "function": {"name": "get_supplier_metrics", "description": "查询供货商指标：涉及订单/分单行数/累计金额/商品 top10。", "parameters": {"type": "object", "properties": {"supplier_id": {"type": "integer"}, "start_date": {"type": "string"}, "end_date": {"type": "string"}}, "required": ["supplier_id"]}}},
    {"type": "function", "function": {"name": "get_client_metrics", "description": "查询客户（学校）采购指标：单数/GMV/按食堂拆分/按品类拆分/配送商分布。", "parameters": {"type": "object", "properties": {"client_id": {"type": "integer"}, "canteen_id": {"type": "integer"}, "start_date": {"type": "string"}, "end_date": {"type": "string"}}, "required": ["client_id"]}}},
    {"type": "function", "function": {"name": "get_factory_metrics", "description": "查询厂家分单指标。", "parameters": {"type": "object", "properties": {"factory_id": {"type": "integer"}, "start_date": {"type": "string"}, "end_date": {"type": "string"}}, "required": ["factory_id"]}}},
    # ===== 类 D：IoT/资产/账务监督 =====
    {"type": "function", "function": {"name": "get_device_status", "description": "查询设备状态清单与最近心跳。可只看离线/只看在线。", "parameters": {"type": "object", "properties": {"device_type": {"type": "string"}, "delivery_id": {"type": "integer"}, "status": {"type": "string"}, "online_only": {"type": "boolean"}, "offline_only": {"type": "boolean"}, "limit": {"type": "integer"}}}}},
    {"type": "function", "function": {"name": "get_iot_recent", "description": "查询近 N 小时（默认 1 小时，上限 168）的 IoT 数据流。", "parameters": {"type": "object", "properties": {"device_type": {"type": "string", "enum": ["scale", "gps", "camera", "sensor"]}, "device_id": {"type": "string"}, "hours": {"type": "integer"}, "limit": {"type": "integer"}}}}},
    {"type": "function", "function": {"name": "get_overdue_bills", "description": "查询超期未付/未确认的账单。scope=payment 取超期未付；scope=confirm 取超期未确认。", "parameters": {"type": "object", "properties": {"scope": {"type": "string", "enum": ["payment", "confirm"]}, "limit": {"type": "integer"}}}}},
    # ===== 类 E：报价/对比 =====
    {"type": "function", "function": {"name": "get_supplier_quotes", "description": "查询供货商报价。", "parameters": {"type": "object", "properties": {"supplier_id": {"type": "integer"}, "product_id": {"type": "integer"}, "limit": {"type": "integer"}}}}},
    {"type": "function", "function": {"name": "compare_supplier_vs_xinfadi", "description": "对比供货商报价与商品参考价、全国农产品近期均价。", "parameters": {"type": "object", "properties": {"product_id": {"type": "integer"}, "supplier_id": {"type": "integer"}}, "required": ["product_id"]}}},
    # ===== 全国农产品价格 / 报告 / schema =====
    {"type": "function", "function": {"name": "get_national_ag_price", "description": "查询全国农产品价格（中农价格网）历史行情。", "parameters": {"type": "object", "properties": {"product_name": {"type": "string"}, "query_text": {"type": "string"}}}}},
    {"type": "function", "function": {"name": "get_national_ag_forecast_price", "description": "查询全国农产品价格未来预测（与数据挖掘中心全国 Tab 一致）。", "parameters": {"type": "object", "properties": {"query_text": {"type": "string"}, "product_name": {"type": "string"}, "days": {"type": "integer"}}}}},
    {"type": "function", "function": {"name": "get_xinfadi_price", "description": "（兼容）同 get_national_ag_price。", "parameters": {"type": "object", "properties": {"product_name": {"type": "string"}, "query_text": {"type": "string"}}}}},
    {"type": "function", "function": {"name": "get_xinfadi_forecast_price", "description": "（兼容）同 get_national_ag_forecast_price。", "parameters": {"type": "object", "properties": {"query_text": {"type": "string"}, "product_name": {"type": "string"}, "days": {"type": "integer"}}}}},
    {"type": "function", "function": {"name": "generate_report", "description": "生成日报/周报/月报/财务报表。when_to_use: 『给我整一份昨天的监管日报』『中农食迅的日报』。when_not_to_use: 单个 KPI 数字或简单概况用 get_kpi_summary，不要为一条数字生成整份报告。昨天/今天必须传 start_date+end_date（ISO）。", "parameters": {"type": "object", "properties": {"report_type": {"type": "string", "enum": ["daily", "weekly", "monthly"], "description": "日报/周报/月报"}, "subject_name": {"type": "string", "description": "主体名字（中文公司名/学校名/食堂名/用户名），AI 内部会自动模糊匹配"}, "subject_type": {"type": "string", "enum": ["client", "delivery", "supplier", "factory", "canteen"]}, "subject_id": {"type": "integer"}, "start_date": {"type": "string"}, "end_date": {"type": "string"}}}}},
    {"type": "function", "function": {"name": "lookup_entity_by_name", "description": "把中文名字（如『中农食迅』『教师食堂』『北京第一实验小学』）模糊匹配到系统实体 ID。返回候选清单，含 type 和 id。用户问『X 的报告/X 的订单』时若不知道 ID，先调此工具。", "parameters": {"type": "object", "properties": {"name": {"type": "string"}, "role": {"type": "string", "enum": ["client", "delivery", "supplier", "factory", "canteen"], "description": "限定搜索范围"}, "limit": {"type": "integer"}}, "required": ["name"]}}},
    {"type": "function", "function": {"name": "get_schema_overview", "description": "说明可查询的数据范围。", "parameters": {"type": "object", "properties": {}}}},
    {"type": "function", "function": {"name": "search_docs", "description": "检索操作手册：怎么用、入口在哪、演示账号、操作步骤。when_to_use: 操作/入口/步骤/账号类问题。when_not_to_use: 实时 GMV/订单/告警数字（用 KPI/search 工具）。", "parameters": {"type": "object", "properties": {"query": {"type": "string"}, "limit": {"type": "integer"}}, "required": ["query"]}}},
    {"type": "function", "function": {"name": "ask_clarification", "description": "当主体/时间/品名无法从问题与会话记忆推断，且默认值会误导监管决策时，向用户提出一个简短澄清问题。不要用于已有足够信息的情况。", "parameters": {"type": "object", "properties": {"question": {"type": "string", "description": "一句中文追问"}}, "required": ["question"]}}},
    # ===== 万能只读查询：结构化工具覆盖不到的任意分析问题，自己写 SQL =====
    {"type": "function", "function": {"name": "run_sql", "description": (
        "对监管数据库执行【只读 SELECT】查询，回答结构化工具覆盖不到的任意分析问题"
        "（如『昨天哪个供货商分的钱最多』『各配送商本月异常率排行』等聚合/排序/多表关联）。"
        "规则：只能写单条 SELECT/WITH；自动加 LIMIT；超时与行数受限；禁止写操作与系统库。"
        "写 SQL 前若不确定表结构，先调 get_schema_overview。日期范围用 MySQL 函数（如 CURDATE()、DATE_SUB）。"
        "when_not_to_use: KPI 概览用 get_kpi_summary，订单明细列表用 search_orders，已有专用工具时勿写 SQL。"
    ), "parameters": {"type": "object", "properties": {"sql": {"type": "string", "description": "单条只读 SELECT 语句"}, "purpose": {"type": "string", "description": "这条查询想回答什么（中文，便于日志与解释）"}}, "required": ["sql"]}}},
]


# ===== 只读 SQL 工具：安全校验 + 执行 + 脱敏 =====
_SQL_FORBIDDEN = re.compile(
    r"\b(insert|update|delete|drop|alter|create|truncate|replace|grant|revoke|rename|"
    r"lock|unlock|call|use|load_file|outfile|dumpfile|benchmark|sleep|"
    r"into\s+outfile|into\s+dumpfile)\b",
    re.I,
)
_SQL_SYS_SCHEMA = re.compile(r"\b(information_schema|mysql|performance_schema|sys)\b", re.I)
_PII_PHONE_COL = ("phone", "mobile", "tel", "电话", "手机", "contact")
_PII_NAME_COL = ("contact_name", "联系人", "linkman", "owner_name", "real_name")
_PII_ADDR_COL = ("address", "addr", "地址")


def _validate_readonly_sql(sql: str) -> tuple[bool, str, str]:
    """校验为单条只读 SELECT/WITH。返回 (ok, 规范化sql 或 错误信息, error)。"""
    s = (sql or "").strip()
    # 去掉结尾分号；中间不允许再出现分号（防多语句）
    if s.endswith(";"):
        s = s[:-1].strip()
    if not s:
        return False, "", "SQL 为空"
    if ";" in s:
        return False, "", "只允许执行单条语句"
    if not re.match(r"^\(?\s*(select|with)\b", s, re.I):
        return False, "", "只允许 SELECT/WITH 查询"
    if _SQL_SYS_SCHEMA.search(s):
        return False, "", "禁止访问系统库（information_schema/mysql/...）"
    # 优先用 sqlglot 做 AST 级校验，失败则回退正则
    try:
        import sqlglot
        from sqlglot import exp

        statements = sqlglot.parse(s, read="mysql")
        if len(statements) != 1 or statements[0] is None:
            return False, "", "只允许单条查询"
        root = statements[0]
        if root.key not in ("select", "union", "with", "subquery"):
            return False, "", "顶层必须是 SELECT/WITH 查询"
        forbidden_nodes = (
            exp.Insert, exp.Update, exp.Delete, exp.Drop, exp.Create, exp.Alter,
            exp.Command, exp.Set, exp.Use, exp.TruncateTable,
        )
        for node in root.walk():
            if isinstance(node, forbidden_nodes):
                return False, "", "检测到非只读操作，已拒绝"
    except ImportError:
        if _SQL_FORBIDDEN.search(s):
            return False, "", "检测到非只读关键字，已拒绝"
    except Exception:
        # 解析异常：用正则兜底，宁可拒绝
        if _SQL_FORBIDDEN.search(s):
            return False, "", "SQL 无法安全解析，已拒绝"
    if _SQL_FORBIDDEN.search(s):
        return False, "", "检测到非只读关键字，已拒绝"
    return True, s, ""


def _mask_sql_cell(col_name: str, value: Any) -> Any:
    if value is None:
        return None
    col = str(col_name or "").lower()
    text_value = str(value)
    if any(h in col for h in _PII_PHONE_COL):
        return _mask_phone(text_value) if re.search(r"\d{7,}", text_value) else text_value
    if any(h in col for h in _PII_NAME_COL):
        return _mask_name(text_value)
    if any(h in col for h in _PII_ADDR_COL):
        return (text_value[:6] + "…") if len(text_value) > 6 else text_value
    # 自由文本里的手机号也兜底掩码
    if isinstance(value, str):
        return _mask_phone(text_value)
    if isinstance(value, Decimal):
        return _num(value)
    if isinstance(value, (datetime, date)):
        return value.isoformat()
    return value


async def _tool_run_sql(args: dict[str, Any]) -> dict[str, Any]:
    raw_sql = str(args.get("sql") or "").strip()
    purpose = str(args.get("purpose") or "").strip()
    ok, norm_sql, err = _validate_readonly_sql(raw_sql)
    if not ok:
        return {"type": "error", "title": "SQL 被拒绝", "rows": [], "columns": [],
                "summary": f"该查询未通过只读安全校验：{err}", "sql": raw_sql, "error": err}
    cap = int(settings.ai_sql_row_limit)
    exec_sql = norm_sql if re.search(r"\blimit\b", norm_sql, re.I) else f"{norm_sql} LIMIT {cap}"
    try:
        async with ReadOnlySession() as ro:
            try:
                await ro.execute(text("SET SESSION MAX_EXECUTION_TIME = :t"), {"t": int(settings.ai_sql_timeout_ms)})
            except Exception:
                pass  # 个别版本不支持则忽略，仍有行数上限兜底
            result = await ro.execute(text(exec_sql))
            col_names = list(result.keys())
            raw_rows = result.fetchmany(cap)
    except Exception as exc:  # noqa: BLE001
        msg = str(exc)
        return {"type": "error", "title": "查询失败", "rows": [], "columns": [],
                "summary": f"SQL 执行失败：{msg[:200]}", "sql": exec_sql, "error": msg[:200]}
    rows = [
        {col: _mask_sql_cell(col, val) for col, val in zip(col_names, row)}
        for row in raw_rows
    ]
    columns = [{"key": c, "label": c} for c in col_names]
    return {
        "type": "table",
        "title": purpose or "查询结果",
        "columns": columns,
        "rows": rows,
        "row_count": len(rows),
        "sql": exec_sql,
        "summary": f"查询返回 {len(rows)} 行" + (f"（已截断至 {cap} 行）" if len(rows) >= cap else "") + "。",
    }


async def _dispatch_tool(name: str, args: dict[str, Any], db: AsyncSession) -> dict[str, Any]:
    # 特殊：写操作拒绝 / 闲聊（不查数据，仅返回静态消息）
    if name == "run_sql":
        return await _tool_run_sql(args)
    if name == "_refuse_write":
        return {
            "type": "notice",
            "title": "只读权限提示",
            "rows": [],
            "summary": (
                "监管端 AI 仅提供只读数据分析，无法执行写操作（取消/删除/修改/审批/标记等）。"
                "如需变更业务数据，请由对应业务端操作：客户端取消订单、运营端审批工单、配送商更新分单等。"
            ),
        }
    if name == "_greeting":
        return {
            "type": "notice",
            "title": "监管端 AI 助手",
            "rows": [],
            "summary": (
                "你好！我是监管端 AI 分析助手，可帮你查询订单、账单、合约、配送、报价、IoT 等任意业务数据，并生成各主体的日/周/月报。"
                "例如：『今天有几单』『中农食迅最近一周履约』『北京第一实验小学的日报』『后天大白菜多少钱』。"
            ),
        }
    # 基础
    if name == "get_kpi_summary":
        return await _tool_kpi_summary(db, args)
    if name == "get_region_rank":
        return await _tool_rank(db, args, "region")
    if name == "get_category_distribution":
        return await _tool_rank(db, args, "category")
    if name == "get_top_goods":
        return await _tool_rank(db, args, "goods")
    if name == "get_top_members":
        return await _tool_rank(db, args, "customer")
    if name == "get_daily_trend":
        return await _tool_daily_trend(db, args)
    if name in ("get_ops_alerts", "search_alerts"):
        return await _tool_search_alerts(db, args) if name == "search_alerts" else await _tool_ops_alerts(db, args)
    if name == "get_today_orders":
        return await _tool_today_orders(db, args)
    # 类 A 实体溯源
    if name == "get_order_detail":
        return await _tool_get_order_detail(db, args)
    if name == "get_statement_detail":
        return await _tool_get_statement_detail(db, args)
    if name == "get_ticket_detail":
        return await _tool_get_ticket_detail(db, args)
    if name == "get_contract_detail":
        return await _tool_get_contract_detail(db, args)
    if name == "get_user_detail":
        return await _tool_get_user_detail(db, args)
    # 类 B 列表搜索
    if name == "search_orders":
        return await _tool_search_orders(db, args)
    if name == "search_statements":
        return await _tool_search_statements(db, args)
    if name == "search_tickets":
        return await _tool_search_tickets(db, args)
    if name == "search_quality_reports":
        return await _tool_search_quality_reports(db, args)
    # 类 C 主体聚合
    if name == "get_delivery_metrics":
        return await _tool_get_delivery_metrics(db, args)
    if name == "get_supplier_metrics":
        return await _tool_get_supplier_metrics(db, args)
    if name == "get_client_metrics":
        return await _tool_get_client_metrics(db, args)
    if name == "get_factory_metrics":
        return await _tool_get_factory_metrics(db, args)
    # 类 D IoT / 资产 / 账务
    if name == "get_device_status":
        return await _tool_get_device_status(db, args)
    if name == "get_iot_recent":
        return await _tool_get_iot_recent(db, args)
    if name == "get_overdue_bills":
        return await _tool_get_overdue_bills(db, args)
    # 类 E 报价对比
    if name == "get_supplier_quotes":
        return await _tool_get_supplier_quotes(db, args)
    if name == "compare_supplier_vs_xinfadi":
        return await _tool_compare_supplier_vs_xinfadi(db, args)
    # 全国农产品价格 / 报告 / schema
    if name in ("get_national_ag_price", "get_xinfadi_price", "get_national_ag_forecast_price", "get_xinfadi_forecast_price"):
        return await _run_national_price_tool(db, name, args)
    if name == "generate_report":
        return await _tool_generate_report(db, args)
    if name == "lookup_entity_by_name":
        return await _tool_lookup_entity_by_name(db, args)
    if name == "get_schema_overview":
        return await _tool_schema_overview(db, args)
    if name == "search_docs":
        return await _tool_search_docs(db, args)
    if name == "ask_clarification":
        question = str(args.get("question") or args.get("query") or "请补充时间范围或主体名称。").strip()
        return {"type": "clarify", "title": "需要澄清", "question": question, "rows": [], "summary": question}
    return {"type": "error", "summary": f"未知工具：{name}", "rows": []}


_REPORT_NOISE_PREFIXES = (
    "给我整一份", "给我整", "给我", "请帮我", "请", "帮我", "我要看", "我要", "我想要", "我想",
    "生成", "做", "出", "看一下", "查一下", "查", "看", "整一份",
)
_REPORT_NOISE_TIME = ("今天", "今日", "这个月", "这周", "这月", "本周", "本月", "上个月", "上周", "上月", "昨天", "昨日", "明天", "近期", "最近")
_REPORT_NOISE_FILLER = ("一个", "一份", "一张", "份", "张", "整", "整份")

# 通用概念词：是"概念"不是实体名，绝不能当作 subject_name / lookup_entity_by_name 的实体去匹配。
_NON_ENTITY_WORDS = (
    "财务", "财报", "经营", "营收", "营业额", "营业", "收入", "支出",
    "统计", "汇总", "数据", "报表", "账目", "资金", "现金流", "盈亏", "成本",
    "整体", "全局", "全部", "系统", "平台", "监管",
)
# 触发"全局财务报表"的财务概念词。
_FINANCIAL_WORDS = ("财务", "财报", "经营", "营收", "营业", "收入", "支出", "资金", "现金流", "盈亏", "账目")

# 领域分类用关键词（强化提示词+规则，无额外 LLM 调用）。
_NATIONAL_PRICE_DOMAIN_KW = (
    "全国农产品", "中农价格网", "中农价格", "农产品价格", "菜价", "批发价", "行情", "新发地",
)
_SYSTEM_DOMAIN_KW = (
    "订单", "下单", "账单", "对账", "结算", "超期", "工单", "投诉", "质检", "质量",
    "设备", "秤", "GPS", "传感器", "离线", "IoT", "iot", "履约", "完成率", "报价",
    "告警", "预警", "异常", "区域", "品类", "GMV", "kpi", "KPI",
    "财务", "财报", "经营", "营收", "营业", "报表", "日报", "周报", "月报", "报告", "简报",
    "客户", "学校", "小学", "中学", "幼儿园", "食堂", "配送商", "物流", "供货商", "厂家",
    "合约", "采购", "今天有几单", "今日订单",
    "生意", "咋样", "概况", "链路", "稽核", "演示", "账号", "密码", "那边", "最近", "查一下",
)


def _classify_domain(text_value: str) -> str:
    """把问题分到三类：national_price、system、out_of_scope。"""
    t = text_value or ""
    product_kw = _extract_ag_product_keyword(t)
    price_hit = any(k in t for k in ("价格", "价钱", "多少钱", "什么价", "参考价", "预测"))
    biz_hit = any(k in t for k in ("订单", "GMV", "账单", "对账", "稽核", "告警", "工单", "履约", "生意"))
    if biz_hit and not (price_hit and product_kw not in ("", t)):
        return "system"
    if any(k in t for k in _NATIONAL_PRICE_DOMAIN_KW):
        return "national_price"
    if price_hit and product_kw not in ("", t):
        return "national_price"
    if any(k in t for k in _SYSTEM_DOMAIN_KW):
        return "system"
    if price_hit:
        return "national_price"
    return "out_of_scope"


def _extract_month_range_from_text(text_value: str) -> Optional[tuple[date, date]]:
    """从文本里提取『X 月/X 月份』范围，自动落到当年。"""
    import calendar
    m = re.search(r"(?<!\d)(1[0-2]|[1-9])\s*月", text_value)
    if not m:
        return None
    month = int(m.group(1))
    year = _today().year
    last_day = calendar.monthrange(year, month)[1]
    start = date(year, month, 1)
    end = date(year, month, last_day)
    today = _today()
    if end > today:
        end = today
    return start, end


def _extract_day_range_from_text(text_value: str) -> Optional[tuple[date, date]]:
    """从文本提取相对单日范围：昨天/今天/前天。"""
    t = (text_value or "").strip()
    if any(k in t for k in ("前天",)):
        d = _today() - timedelta(days=2)
        return d, d
    if any(k in t for k in ("昨天", "昨日")):
        d = _today() - timedelta(days=1)
        return d, d
    if any(k in t for k in ("今天", "今日", "当日", "当天")):
        d = _today()
        return d, d
    return None


def _apply_day_range_to_report_args(text_value: str, gen_args: dict[str, Any]) -> dict[str, Any]:
    day_range = _extract_day_range_from_text(text_value)
    if day_range:
        gen_args = {**gen_args, "start_date": day_range[0].isoformat(), "end_date": day_range[1].isoformat()}
    return gen_args


def _extract_subject_name(text_value: str) -> Optional[str]:
    """从含"X 的 日报/周报/月报"的句子里提取主体名。
    例：『给我生成一个今天中农食迅的日报』→ 『中农食迅』。"""
    for tok in ("月报", "周报", "日报", "报告", "简报", "报表"):
        if tok not in text_value:
            continue
        head = text_value.split(tok)[0]
        # 包含式清除时间词与填充词（覆盖前缀/中缀/后缀，如"中农食迅这个月的"）
        for t in _REPORT_NOISE_TIME:
            head = head.replace(t, "")
        for s in _REPORT_NOISE_FILLER:
            head = head.replace(s, "")
        head = head.replace("的", "").strip()
        # 反复剥前缀动词
        changed = True
        while changed:
            changed = False
            for p in _REPORT_NOISE_PREFIXES:
                if head.startswith(p):
                    head = head[len(p):].strip()
                    changed = True
        head = head.strip(" ,，:：")
        # 剥掉通用概念词（财务/经营/统计/汇总…），剩下的才是真正的主体名；纯概念词 → None
        for w in _NON_ENTITY_WORDS:
            head = head.replace(w, "")
        head = head.strip(" ,，:：")
        if head:
            return head
        return None
    return None


_WRITE_ACTION_KW = (
    "取消", "删除", "删了", "删掉", "删一下", "去掉这", "改成", "改为", "修改成", "改一下", "调一下",
    "强制结算", "强制确认", "强制关闭", "审批通过", "审批拒绝", "标记为", "标记成", "改账", "调账",
    "下架", "上架", "重新分单", "帮我改", "帮我删", "帮我取消", "新增", "添加一条", "创建订单", "录入",
)
_GREETING_KW = ("你好", "您好", "hi", "Hi", "hello", "Hello", "嗨", "在吗", "在不在", "在么", "在?")
_VAGUE_ORDER_PATTERN = re.compile(r"^查(一下|下)?\s*订单\s*$")


def _heuristic_tool(text_value: str, entities: Optional[dict[str, Any]] = None) -> tuple[str, dict[str, Any]]:
    entities = entities or {}
    t = (text_value or "").strip()
    # 写操作明确拒绝
    if any(k in t for k in _WRITE_ACTION_KW):
        return "_refuse_write", {"text": t}
    # 闲聊
    if t in _GREETING_KW or t.lower() in {k.lower() for k in _GREETING_KW}:
        return "_greeting", {}
    fast_price = _fast_price_query_from_text(text_value)
    if fast_price:
        return fast_price
    # 文本中或上下文里含订单号 + 详情类问法 → 实体溯源
    explicit_order_no = None
    for pat in _ORDER_NO_PATTERNS:
        m = pat.search(text_value)
        if m:
            explicit_order_no = m.group(0)
            break
    if explicit_order_no:
        return "get_order_detail", {"order_no": explicit_order_no}
    if "是谁" in text_value or text_value.rstrip("?？").endswith("是什么"):
        who_name = re.sub(r"(是谁|是什么|哪位).*$", "", text_value).strip()
        for noise in ("请", "帮我", "查一下", "查", "看一下", "看", "我想知道"):
            who_name = who_name.replace(noise, "")
        who_name = who_name.strip(" ,，:：")
        if len(who_name) >= 2:
            return "lookup_entity_by_name", {"name": who_name, "limit": 5}
    has_detail_kw = any(k in text_value for k in ("详情", "信息", "是谁", "谁下", "下的", "明细"))
    if has_detail_kw and entities.get("latest_order_no"):
        return "get_order_detail", {"order_no": entities["latest_order_no"]}
    if has_detail_kw and entities.get("latest_statement_no"):
        return "get_statement_detail", {"statement_no": entities["latest_statement_no"]}
    # 用户提问含代词（"这个单/该订单"）且上下文有订单号 → 回到 get_order_detail（含异常、工单、分单等完整信息）
    if entities.get("has_pronoun") and entities.get("latest_order_no"):
        return "get_order_detail", {"order_no": entities["latest_order_no"]}
    if entities.get("has_pronoun") and entities.get("latest_statement_no"):
        return "get_statement_detail", {"statement_no": entities["latest_statement_no"]}
    # 主体类业务问句：「X 食堂今天买了多少」「X 学校的订单」「中农食迅履约」等 → lookup + 指标
    # 必须放在新发地行情/价格匹配之前，避免被"多少钱"误导
    _entity_noun_kw = ("食堂", "学校", "小学", "中学", "幼儿园", "客户", "配送商", "物流", "供货商", "厂家", "中农", "食迅")
    _business_action_kw = ("买", "采购", "下单", "订单", "GMV", "金额", "履约", "完成率", "异常率", "多少钱", "多少", "营业", "销售", "总额")
    if any(n in text_value for n in _entity_noun_kw) and any(a in text_value for a in _business_action_kw):
        cleaned = text_value
        for noise in ("给我", "请", "帮我", "查一下", "查", "看一下", "看", "我要", "我想", "?", "？", "怎么样", "怎样", "如何", "情况"):
            cleaned = cleaned.replace(noise, "")
        for t in ("这个月", "这周", "本周", "本月", "上个月", "上周", "上月", "今天", "今日", "昨天", "近期", "最近", "买了", "采购了", "下了", "的"):
            cleaned = cleaned.replace(t, "")
        for a in _business_action_kw:
            cleaned = cleaned.replace(a, "")
        # 剥掉业务动作词残留的悬挂字（如"履约率"去掉"履约"后残留的"率"）
        cleaned = cleaned.replace("履约率", "").replace("完成率", "").replace("率", "").replace("额", "")
        # 月份词清洗
        cleaned = re.sub(r"(1[0-2]|[1-9])\s*月份?", "", cleaned)
        cleaned = cleaned.strip()
        if cleaned:
            return "lookup_entity_by_name", {"name": cleaned, "limit": 5}
    # 账单/对账
    if any(k in text_value for k in ("超期未付", "超期未结", "逾期账单", "超期账单")):
        return "get_overdue_bills", {"scope": "payment", "limit": 20}
    if any(k in text_value for k in ("超期未确认", "待确认账单")):
        return "get_overdue_bills", {"scope": "confirm", "limit": 20}
    if any(k in text_value for k in ("账单", "对账", "结算")):
        return "search_statements", {"limit": 20}
    # 工单
    if any(k in text_value for k in ("工单", "投诉", "异常订单")):
        return "search_tickets", {"limit": 20}
    # 质检
    if any(k in text_value for k in ("质检", "质量报告", "化验")):
        return "search_quality_reports", {"limit": 20}
    # 设备 IoT
    if any(k in text_value for k in ("设备", "秤", "GPS", "传感器", "离线", "在线")):
        if "离线" in text_value:
            return "get_device_status", {"offline_only": True, "limit": 30}
        return "get_device_status", {"limit": 30}
    if any(k in text_value for k in ("IoT", "iot", "温度", "湿度", "上报数据")):
        return "get_iot_recent", {"hours": 1, "limit": 30}
    # 履约/主体指标
    if "履约" in text_value or "完成率" in text_value:
        return "get_kpi_summary", {"scope": "range", "start_date": (_today() - timedelta(days=6)).isoformat(), "end_date": _today().isoformat()}
    # 报价对比
    if any(k in text_value for k in ("报价", "对比新发地", "比新发地")):
        return "get_supplier_quotes", {"limit": 20}
    # schema / 报告
    if any(k in text_value for k in ("能查", "哪些数据", "表结构", "数据范围")):
        return "get_schema_overview", {}
    # 财务/经营报表（无主体）→ 全局财务汇总。财务/经营是概念，不能当主体名。
    if any(k in text_value for k in _FINANCIAL_WORDS):
        fin_args: dict[str, Any] = {"report_type": "financial"}
        mr = _extract_month_range_from_text(text_value)
        if mr:
            fin_args["start_date"] = mr[0].isoformat()
            fin_args["end_date"] = mr[1].isoformat()
        return "generate_report", fin_args
    if any(k in text_value for k in ("日报", "周报", "月报", "报告", "简报", "报表", "导出")):
        if any(k in text_value for k in ("月报", "本月", "这个月", "上个月", "上月", "当月")):
            rt = "monthly"
        elif any(k in text_value for k in ("周报", "本周", "这周", "上周")):
            rt = "weekly"
        else:
            rt = "daily"
        gen_args: dict[str, Any] = {"report_type": rt}
        subject_name = _extract_subject_name(text_value)
        if subject_name:
            gen_args["subject_name"] = subject_name
        gen_args = _apply_day_range_to_report_args(text_value, gen_args)
        return "generate_report", gen_args
    # 全国农产品价格
    np_tool = _national_price_tool_from_text(text_value)
    if np_tool:
        return np_tool
    # 趋势（业务 KPI，非菜价）
    if any(k in text_value for k in ("趋势", "折线", "走势")) and not any(
        k in text_value for k in ("菜价", "农产品", "批发", "多少钱", "价格预测")
    ):
        return "get_daily_trend", {"start_date": (_today() - timedelta(days=6)).isoformat(), "end_date": _today().isoformat()}
    # 告警
    if any(k in text_value for k in ("异常", "预警", "告警")):
        return "search_alerts", {"limit": 20}
    # 订单列表
    if "今天" in text_value and any(k in text_value for k in ("单", "订单")):
        return "get_today_orders", {}
    # 模糊订单查询 → 默认今天（回复中说明假设）
    if re.match(r"^查(一下|下)?\s*订单\s*$", t) or t in ("查订单", "看看订单", "订单呢"):
        d = _today().isoformat()
        return "search_orders", {"start_date": d, "end_date": d, "limit": 20}
    if any(k in text_value for k in ("订单", "下单")):
        month_range = _extract_month_range_from_text(text_value)
        if month_range:
            return "search_orders", {"start_date": month_range[0].isoformat(), "end_date": month_range[1].isoformat(), "limit": 50}
        return "search_orders", {"limit": 20}
    # 品类分布 / 饼图（优先于普通排行）
    if any(k in text_value for k in ("饼图", "占比", "比例", "分布")) and any(
        k in text_value for k in ("品类", "分类", "一级分类", "一级")
    ):
        today_only = any(k in text_value for k in ("今天", "今日", "当日", "当天"))
        return "get_category_distribution", {
            "start_date": _today().isoformat() if today_only else (_today() - timedelta(days=29)).isoformat(),
            "end_date": _today().isoformat(),
            "limit": 10,
            "chart_type": "pie",
        }
    # 排行
    if any(k in text_value for k in ("区域", "哪个区")):
        return "get_region_rank", {"start_date": (_today() - timedelta(days=29)).isoformat(), "end_date": _today().isoformat(), "limit": 10}
    if "商品" in text_value:
        return "get_top_goods", {"start_date": (_today() - timedelta(days=29)).isoformat(), "end_date": _today().isoformat(), "limit": 10}
    if "品类" in text_value or "分类" in text_value:
        today_only = any(k in text_value for k in ("今天", "今日", "当日", "当天"))
        return "get_category_distribution", {
            "start_date": _today().isoformat() if today_only else (_today() - timedelta(days=29)).isoformat(),
            "end_date": _today().isoformat(),
            "limit": 10,
            "chart_type": "pie" if "饼图" in text_value else "bar",
        }
    # 会话主体续问（那边配送商最近咋样 — 优先于泛化「咋样→KPI」）
    session_name = str(entities.get("session_subject_name") or "").strip()
    if session_name and any(
        k in text_value for k in ("那边", "这家", "该", "最近咋样", "最近怎么样", "最近如何", "最近咋", "咋样", "怎么样")
    ):
        return "lookup_entity_by_name", {"name": session_name, "limit": 5}
    # 口语化经营/KPI（昨天生意咋样、最近概况）
    if any(k in text_value for k in ("生意", "咋样", "概况", "经营情况", "卖得怎么样", "卖得咋样")):
        day_range = _extract_day_range_from_text(text_value)
        if day_range:
            return "get_kpi_summary", {
                "scope": "range",
                "start_date": day_range[0].isoformat(),
                "end_date": day_range[1].isoformat(),
            }
        if any(k in text_value for k in ("最近", "近期")):
            return "get_kpi_summary", {
                "scope": "range",
                "start_date": (_today() - timedelta(days=6)).isoformat(),
                "end_date": _today().isoformat(),
            }
        return "get_kpi_summary", {"scope": "today"}
    # 手册/入口（链路在哪查 — 避免落入 out_of_scope）
    if any(k in text_value for k in ("链路", "手册", "入口", "浮窗", "AI助手", "AI 助手")) and any(
        k in text_value for k in ("在哪", "哪里", "怎么", "如何", "查")
    ):
        return "search_docs", {"query": text_value or "操作手册"}
    # KPI 兜底：识别时间范围（这个月/本月/X月/最近N天/本周）
    month_range = _extract_month_range_from_text(text_value)
    if month_range:
        return "get_kpi_summary", {"scope": "range", "start_date": month_range[0].isoformat(), "end_date": month_range[1].isoformat()}
    if any(k in text_value for k in ("这个月", "本月", "当月")):
        return "get_kpi_summary", {"scope": "range", "start_date": _today().replace(day=1).isoformat(), "end_date": _today().isoformat()}
    if any(k in text_value for k in ("这周", "本周")):
        return "get_kpi_summary", {"scope": "range", "start_date": (_today() - timedelta(days=_today().weekday())).isoformat(), "end_date": _today().isoformat()}
    m_days = re.search(r"最近\s*(\d{1,3})\s*天", text_value)
    if m_days:
        n = max(1, min(365, int(m_days.group(1))))
        return "get_kpi_summary", {"scope": "range", "start_date": (_today() - timedelta(days=n - 1)).isoformat(), "end_date": _today().isoformat()}
    # 既不属于本系统、也不属于新发地行情/预测 → 告知可问范围，不死胡同
    if _classify_domain(text_value) == "out_of_scope":
        return "get_schema_overview", {}
    return "get_kpi_summary", {"scope": "today"}


def _system_prompt() -> str:
    return build_system_prompt(_today().isoformat())


def _context_reference_prompt(entities: dict[str, Any]) -> Optional[str]:
    if not entities:
        return None
    parts: list[str] = []
    if entities.get("latest_order_no"):
        parts.append(f"最新订单号: {entities['latest_order_no']}")
    if entities.get("latest_statement_no"):
        parts.append(f"最新账单号: {entities['latest_statement_no']}")
    if not parts:
        return None
    hint = ""
    if entities.get("has_pronoun"):
        hint = "  ⚠️ 用户本轮提问包含“这个单/该订单/这张账单”等代词，请使用上述编号作为工具参数。"
    return "【上下文引用】" + "；".join(parts) + "。" + hint


async def _dashscope_chat(messages: list[dict[str, Any]], tools: Optional[list[dict[str, Any]]] = None) -> dict[str, Any]:
    api_key = (settings.ai_api_key or "").strip()
    if not api_key:
        raise RuntimeError("AI_API_KEY 未配置")
    payload: dict[str, Any] = {
        "model": settings.ai_model_answer or "qwen-plus",
        "messages": messages,
        "temperature": 0.2,
    }
    if tools:
        payload["tools"] = tools
        payload["tool_choice"] = "auto"
    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    url = f"{(settings.ai_base_url or '').rstrip('/')}/chat/completions"
    async with httpx.AsyncClient(timeout=httpx.Timeout(25.0, connect=5.0), trust_env=False) as client:
        response = await client.post(url, headers=headers, json=payload)
        response.raise_for_status()
        return response.json()


def _first_choice(payload: dict[str, Any]) -> dict[str, Any]:
    choices = payload.get("choices") or []
    return choices[0].get("message") if choices else {}


def _artifact_from_tool(tool_result: dict[str, Any]) -> tuple[Optional[dict[str, Any]], str]:
    return artifact_from_tool(tool_result)


def _direct_price_reply(tool_result: dict[str, Any]) -> tuple[str, Optional[dict[str, Any]], str]:
    if not isinstance(tool_result, dict):
        return "价格预测结果暂不可用，请稍后重试。", None, ""
    if not tool_result.get("ok"):
        status = str(tool_result.get("status") or "")
        product = str(tool_result.get("product") or tool_result.get("query") or "该品名")
        if status == "ambiguous":
            candidates = tool_result.get("candidates") or []
            text_value = "、".join(str(x) for x in candidates[:8]) if isinstance(candidates, list) and candidates else "无候选"
            return f"我找到多个可能品名：{text_value}。请你回复更准确的品名后我再给明天/未来价格。", None, ""
        if status == "not_found":
            return f"当前库里没有“{product}”这个品名的预测数据。请换一个更准确的品名试试。", None, ""
        return str(tool_result.get("summary") or tool_result.get("message") or f"{product} 当前预测暂不可用。"), None, ""

    product = str(tool_result.get("product") or "该品名")
    mode = str(tool_result.get("mode") or "future")
    rows = tool_result.get("future") or tool_result.get("rows") or []
    if not isinstance(rows, list) or not rows:
        return f"{product} 当前暂无可用未来价格序列，请先训练或刷新模型。", None, ""
    if mode == "tomorrow":
        first = rows[0]
        price = _num(first.get("yhat", first.get("price")))
        confidence = _num(first.get("confidence"))
        trend = {"up": "上涨", "down": "下行", "flat": "持平", "stable": "持平"}.get(str(first.get("trend") or "flat"), "持平")
        day_label = "后日" if int(tool_result.get("target_offset") or 1) == 2 else "明日"
        card = {
            "type": "kpi",
            "title": f"{day_label}价格预测",
            "kpis": [
                {"label": "品名", "value": product},
                {"label": f"{day_label}参考价", "value": f"{price:.3f} 元/斤"},
                {"label": "置信度", "value": f"{confidence * 100:.1f}%"},
                {"label": "趋势", "value": trend},
            ],
            "rows": [],
        }
        src = str(tool_result.get("source") or "")
        src_note = (
            "全国农产品价格预测快照"
            if src == "zgncpjgw_snapshot"
            else "全国行情实时锚定"
        )
        unified = "（已统一按全国农产品价格口径，不再使用新发地数据）" if tool_result.get("unified_from_xinfadi") else ""
        reply = (
            f"{product} 的{day_label}参考价（{first.get('date') or '—'}）为 **{price:.3f} 元/斤**，"
            f"置信度 {confidence * 100:.1f}%，趋势{trend}。该结果来自{src_note}"
            f"（与数据挖掘中心全国农产品价格一致）{unified}"
        )
        return reply, card, ""

    picked = rows[: max(1, min(7, len(rows)))]
    table_rows = []
    for row in picked:
        price = _num(row.get("yhat", row.get("price")))
        table_rows.append(
            {
                "date": str(row.get("date") or "—"),
                "price": round(price, 3),
                "trend": {"up": "上涨", "down": "下行", "flat": "持平", "stable": "持平"}.get(str(row.get("trend") or "flat"), "持平"),
                "confidence": round(_num(row.get("confidence")) * 100, 1),
            }
        )
    card = {
        "type": "trend",
        "title": f"{product} 未来{len(table_rows)}天预测",
        "columns": [
            {"key": "date", "label": "日期"},
            {"key": "price", "label": "预测价(元/斤)"},
            {"key": "trend", "label": "趋势"},
            {"key": "confidence", "label": "置信度(%)"},
        ],
        "kpis": [
            {"label": "品名", "value": product},
            {"label": "模型版本", "value": str(tool_result.get("model_version") or "—")},
            {"label": "首日参考", "value": f"{table_rows[0]['price']:.3f} 元/斤" if table_rows else "—"},
        ],
        "rows": table_rows,
    }
    src = str(tool_result.get("source") or "")
    src_note = "全国农产品价格预测快照" if src == "zgncpjgw_snapshot" else "全国行情实时锚定"
    unified = "（已统一按全国农产品价格口径）" if tool_result.get("unified_from_xinfadi") else ""
    reply = (
        f"已按{src_note}计算 {product} 未来 {len(table_rows)} 天参考价"
        f"（与数据挖掘中心全国农产品价格一致）{unified}"
    )
    return reply, card, ""


def _direct_history_reply(tool_result: dict[str, Any]) -> tuple[str, Optional[dict[str, Any]], str]:
    if not isinstance(tool_result, dict):
        return "历史价格查询暂不可用，请稍后重试。", None, ""
    if tool_result.get("candidates"):
        labels = "、".join(str(x) for x in (tool_result.get("candidates") or [])[:8])
        return f"我找到多个可能品名：{labels}。请补充更准确的规格。", None, ""
    if not tool_result.get("ok"):
        return str(tool_result.get("summary") or "未查到该品名的全国农产品价格。"), None, ""

    product = str(tool_result.get("product") or "该品名")
    target_date = str(tool_result.get("target_date") or "")
    if target_date and tool_result.get("price") is not None:
        price = _num(tool_result.get("price"))
        card = {
            "type": "kpi",
            "title": f"{target_date} 全国均价",
            "kpis": [
                {"label": "品名", "value": product},
                {"label": "日期", "value": target_date},
                {"label": "均价", "value": f"{price:.3f} 元/斤"},
            ],
            "rows": tool_result.get("rows") or [],
        }
        reply = f"{product} 在 **{target_date}** 的全国农产品均价为 **{price:.3f} 元/斤**（数据来源：全国农产品价格库）。"
        return reply, card, ""

    rows = tool_result.get("rows") or []
    card = {
        "type": "national_price",
        "title": tool_result.get("title") or f"{product} 全国行情",
        "columns": tool_result.get("columns") or [],
        "rows": rows[:20],
    }
    reply = str(tool_result.get("summary") or f"{product} 全国农产品历史行情已列出。")
    return reply, card, ""


def _reply_from_national_tool(name: str, tool_result: dict[str, Any]) -> tuple[str, Optional[dict[str, Any]], str]:
    if name in ("get_national_ag_forecast_price", "get_xinfadi_forecast_price"):
        return _direct_price_reply(tool_result)
    return _direct_history_reply(tool_result)


async def _try_question_mode_route(
    mode: str,
    text_value: str,
    history: list[ChatMessage],
    session_id: Optional[str],
    db: AsyncSession,
) -> Optional[dict[str, Any]]:
    """用户手动选择问题类型时，跳过自动意图识别。"""
    m = (mode or "auto").strip().lower()
    if m in ("", "auto"):
        return None

    sid = session_id or f"chat-{datetime.utcnow().strftime('%Y%m%d')}"

    if m == "how_to":
        result = await _dispatch_tool("search_docs", {"query": text_value, "limit": 5}, db)
        data_card, report = _artifact_from_tool(result)
        reply = str(result.get("summary") or "已检索操作手册。")
        _remember_turn(session_id, text_value, reply)
        return {
            "reply": reply,
            "data_card": data_card,
            "report_content": report,
            "session_id": sid,
            "debug": {
                "provider": "question_mode",
                "route": "mode_how_to",
                "intent": "how_to",
                "question_mode": m,
                "tool_calls": [{"name": "search_docs", "arguments": {"query": text_value}, "result_preview": _preview(result)}],
            },
        }

    if m == "report":
        gen_args: dict[str, Any] = {"report_type": "daily"}
        if any(k in text_value for k in ("周报", "本周", "这周")):
            gen_args["report_type"] = "weekly"
        elif any(k in text_value for k in ("月报", "本月", "这个月")):
            gen_args["report_type"] = "monthly"
        gen_args = _apply_day_range_to_report_args(text_value, gen_args)
        subject_name = _extract_subject_name(text_value)
        if subject_name:
            gen_args["subject_name"] = subject_name
        result = await _dispatch_tool("generate_report", gen_args, db)
        data_card, report = _artifact_from_tool(result)
        reply = str(result.get("summary") or report or "报告已生成。")
        _remember_turn(session_id, text_value, reply)
        return {
            "reply": reply,
            "data_card": data_card,
            "report_content": report,
            "session_id": sid,
            "debug": {
                "provider": "question_mode",
                "route": "mode_report",
                "intent": "report",
                "question_mode": m,
                "tool_calls": [{"name": "generate_report", "arguments": gen_args, "result_preview": _preview(result)}],
            },
        }

    if m == "national_price":
        tool_pair = (
            _fast_price_query_from_text(text_value)
            or _fast_price_followup(text_value, history, session_id)
            or _fast_historical_price_from_text(text_value)
            or _national_price_tool_from_text(text_value)
        )
        if not tool_pair:
            return {
                "reply": "请在问题里写明品名，例如「大白菜明天多少钱」或「土豆价格」。",
                "data_card": None,
                "report_content": "",
                "session_id": sid,
                "debug": {"provider": "question_mode", "route": "mode_national_need_product", "intent": "national_price", "question_mode": m, "tool_calls": []},
            }
        name, args = tool_pair
        result = await _run_national_price_tool(db, name, args)
        reply, data_card, report = _reply_from_national_tool(name, result)
        _remember_turn(session_id, text_value, reply)
        return {
            "reply": reply,
            "data_card": data_card,
            "report_content": report,
            "session_id": sid,
            "debug": {
                "provider": "question_mode",
                "route": "mode_national_price",
                "intent": "national_price",
                "question_mode": m,
                "tool_calls": [{"name": name, "arguments": args, "result_preview": _preview(result)}],
            },
        }

    if m == "system_data":
        entities = apply_memory_entities(load_session(session_id), _extract_entities_from_history(history, text_value))
        tool, args = _heuristic_tool(text_value, entities)
        if tool in ("get_national_ag_forecast_price", "get_xinfadi_forecast_price", "get_national_ag_price", "get_xinfadi_price"):
            tool, args = "get_kpi_summary", {"scope": "today"}
        result = await _dispatch_tool(tool, args, db)
        data_card, report = _artifact_from_tool(result)
        reply = str(result.get("summary") or "")
        _remember_turn(session_id, text_value, reply)
        return {
            "reply": reply,
            "data_card": data_card,
            "report_content": report,
            "session_id": sid,
            "debug": {
                "provider": "question_mode",
                "route": "mode_system_data",
                "intent": "analytics",
                "question_mode": m,
                "tool_calls": [{"name": tool, "arguments": args, "result_preview": _preview(result)}],
            },
        }

    return None


async def _local_answer(text_value: str, db: AsyncSession, entities: Optional[dict[str, Any]] = None) -> dict[str, Any]:
    async def _dispatch(name: str, args: dict[str, Any]) -> dict[str, Any]:
        return await _dispatch_tool(name, args, db)

    debug_calls, final, reply = await run_rule_chain(
        text_value,
        heuristic_tool=_heuristic_tool,
        dispatch_tool=_dispatch,
        entities=entities,
        apply_report_day_range=_apply_day_range_to_report_args,
        today=_today(),
    )
    if not debug_calls:
        tool, args = _heuristic_tool(text_value, entities)
        final = await _dispatch_tool(tool, args, db)
        debug_calls = [{"name": tool, "arguments": args, "result_preview": _preview(final)}]
        if tool in ("get_national_ag_forecast_price", "get_xinfadi_forecast_price", "get_national_ag_price", "get_xinfadi_price"):
            reply, data_card, report = _reply_from_national_tool(tool, final)
            return {
                "reply": reply,
                "data_card": data_card,
                "report_content": report,
                "session_id": f"chat-{datetime.utcnow().strftime('%Y%m%d')}",
                "debug": {
                    "provider": "local_rule_fallback",
                    "route": "fallback",
                    "intent": _classify_domain(text_value),
                    "tool_calls": debug_calls,
                },
            }
        reply = str(final.get("summary") or "已完成分析。")
    data_card, report = _artifact_from_tool(final)
    return {
        "reply": reply,
        "data_card": data_card,
        "report_content": report,
        "session_id": f"chat-{datetime.utcnow().strftime('%Y%m%d')}",
        "debug": {
            "provider": "local_rule_fallback",
            "route": "fallback",
            "intent": _classify_domain(text_value),
            "tool_calls": debug_calls,
        },
    }


async def _run_llm_agent(
    text_value: str,
    history: list[ChatMessage],
    db: AsyncSession,
    session_id: Optional[str],
    chat_intent: str,
) -> dict[str, Any]:
    mem = load_session(session_id)
    entities = apply_memory_entities(mem, _extract_entities_from_history(history, text_value))
    base_messages: list[dict[str, Any]] = [{"role": "system", "content": _system_prompt()}]
    mem_hint = mem.to_prompt()
    if mem_hint:
        base_messages.append({"role": "system", "content": mem_hint})
    ctx_hint = _context_reference_prompt(entities)
    if ctx_hint:
        base_messages.append({"role": "system", "content": ctx_hint})
    rel_day = resolve_relative_day(text_value, _today())
    if rel_day and any(k in text_value for k in ("换", "改", "呢", "再来")):
        hint = f"【日期续问】用户希望改用日期区间 {rel_day[0]}~{rel_day[1]}"
        if mem.last_subject_name:
            sid = f"{mem.last_subject_type}#{mem.last_subject_id}" if mem.last_subject_id else (mem.last_subject_type or "")
            hint += (
                f"。必须沿用主体「{mem.last_subject_name}」"
                + (f"（{sid}）" if sid else "")
                + f"，调用 generate_report 或上次同类工具，参数 start_date={rel_day[0]}, end_date={rel_day[1]}"
                + (f", subject_name={mem.last_subject_name}" if mem.last_subject_name else "")
                + "。"
            )
        else:
            hint += "，请沿用会话记忆中的主体/工具类型。"
        base_messages.append({"role": "system", "content": hint})
    hist_limit = int(settings.ai_chat_history_limit)
    for msg in history[-hist_limit:]:
        if msg.role in {"user", "assistant"} and msg.content:
            base_messages.append({"role": msg.role, "content": msg.content})
    if not any(m.get("role") == "user" and m.get("content") == text_value for m in base_messages):
        base_messages.append({"role": "user", "content": text_value})

    tool_calls_debug: list[dict[str, Any]] = []
    messages = list(base_messages)
    max_steps = int(settings.ai_agent_max_steps)
    final_tool_result: dict[str, Any] = {}
    reply = ""
    clarify_question = ""

    for step in range(max_steps):
        resp = await _dashscope_chat(messages, TOOLS)
        message = _first_choice(resp) or {}
        tool_calls = message.get("tool_calls") or []
        if not tool_calls:
            reply = (message.get("content") or "").strip()
            if not reply and step == 0:
                name, args = _heuristic_tool(text_value, entities)
                final_tool_result = await _dispatch_tool(name, args, db)
                update_session_from_tool(session_id, name, args, final_tool_result)
                tool_calls_debug.append({"name": name, "arguments": args, "result_preview": _preview(final_tool_result)})
                reply = str(final_tool_result.get("summary") or "已完成分析。")
            break
        messages.append({"role": "assistant", "content": message.get("content") or "", "tool_calls": tool_calls})
        for call in tool_calls[:6]:
            fn = call.get("function") or {}
            name = fn.get("name") or ""
            try:
                raw_args = fn.get("arguments") or "{}"
                args = raw_args if isinstance(raw_args, dict) else json.loads(raw_args)
            except json.JSONDecodeError:
                args = {}
            tool_result = await _dispatch_tool(name, args, db)
            update_session_from_tool(session_id, name, args, tool_result)
            if tool_result.get("type") != "error" or not final_tool_result:
                final_tool_result = tool_result
            if tool_result.get("type") == "clarify":
                clarify_question = str(tool_result.get("question") or tool_result.get("summary") or "")
            tool_calls_debug.append({
                "name": name,
                "arguments": args,
                "sql": tool_result.get("sql"),
                "result_preview": _preview(tool_result),
                **({"citations": tool_result.get("citations") or []} if name == "search_docs" else {}),
            })
            messages.append({
                "role": "tool",
                "tool_call_id": call.get("id") or f"tool-{len(tool_calls_debug)}",
                "content": json.dumps(compact_tool_result_for_llm(tool_result), ensure_ascii=False),
            })

    draft = reply or clarify_question or str(final_tool_result.get("summary") or "已完成分析。")
    synthesis_used = False
    if (
        (settings.ai_api_key or "").strip()
        and settings.ai_chat_synthesis_enabled
        and should_synthesize(
            tool_calls_debug=tool_calls_debug,
            final_tool_result=final_tool_result,
            reply=reply,
            min_tools=int(settings.ai_chat_synthesis_min_tools),
        )
    ):
        synthesized = await synthesize_answer(
            text_value,
            tool_calls_debug=tool_calls_debug,
            llm_chat=_dashscope_chat,
            first_choice=_first_choice,
            session_hint=mem_hint,
            draft_reply=draft,
            intent=chat_intent,
        )
        if synthesized:
            reply = synthesized
            synthesis_used = True
    if not reply:
        if not (settings.ai_api_key or "").strip():
            raise RuntimeError("AI_API_KEY 未配置")
        final = await _dashscope_chat(
            [*messages, {"role": "system", "content": "请基于以上工具结果用自然、有洞察的中文回答用户，不要输出原始 JSON。"}],
            None,
        )
        reply = (_first_choice(final).get("content") or draft).strip()
    data_card, report = _artifact_from_tool(final_tool_result)
    _remember_turn(session_id, text_value, reply)
    return {
        "reply": reply,
        "data_card": data_card,
        "report_content": report,
        "session_id": session_id or f"chat-{datetime.utcnow().strftime('%Y%m%d')}",
        "debug": {
            "provider": "dashscope",
            "route": "llm_first",
            "intent": chat_intent,
            "domain": _classify_domain(text_value),
            "tool_calls": tool_calls_debug,
            "model": settings.ai_model_answer,
            "steps": len(tool_calls_debug),
            "synthesis": synthesis_used,
            "session_memory": mem.to_debug(),
        },
    }


async def _dispatch(
    text_value: str,
    history: list[ChatMessage],
    db: AsyncSession,
    session_id: Optional[str] = None,
    question_mode: str = "auto",
) -> dict[str, Any]:
    import time

    t0 = time.perf_counter()

    def _finish(result: dict[str, Any]) -> dict[str, Any]:
        ms = int((time.perf_counter() - t0) * 1000)
        dbg = result.setdefault("debug", {})
        if isinstance(dbg, dict):
            dbg["duration_ms"] = ms
        return result

    if any(k in text_value for k in ("北京时间", "几月几号", "今天几号", "现在日期")):
        reply = f"现在按北京时间是 {_today().isoformat()}。"
        _remember_turn(session_id, text_value, reply)
        log_turn(session_id=session_id, user_text=text_value, intent="general", provider="local_rule", route="local_date")
        return _finish({
            "reply": reply,
            "data_card": None,
            "report_content": "",
            "session_id": session_id or f"chat-{datetime.utcnow().strftime('%Y%m%d')}",
            "debug": {"provider": "local_rule", "route": "local_date", "intent": "general", "tool_calls": []},
        })

    safety_tool, _ = _heuristic_tool(text_value, {})
    if safety_tool == "_refuse_write":
        result = await _dispatch_tool("_refuse_write", {}, db)
        reply = str(result.get("summary") or "")
        _remember_turn(session_id, text_value, reply)
        log_turn(session_id=session_id, user_text=text_value, intent="general", provider="local_rule", route="safety")
        return _finish({
            "reply": reply,
            "data_card": None,
            "report_content": "",
            "session_id": session_id or f"chat-{datetime.utcnow().strftime('%Y%m%d')}",
            "debug": {"provider": "local_rule", "route": "safety", "intent": "general", "tool_calls": []},
        })
    if safety_tool == "_greeting":
        result = await _dispatch_tool("_greeting", {}, db)
        reply = str(result.get("summary") or "")
        _remember_turn(session_id, text_value, reply)
        log_turn(session_id=session_id, user_text=text_value, intent="general", provider="local_rule", route="greeting")
        return _finish({
            "reply": reply,
            "data_card": None,
            "report_content": "",
            "session_id": session_id or f"chat-{datetime.utcnow().strftime('%Y%m%d')}",
            "debug": {"provider": "local_rule", "route": "greeting", "intent": "general", "tool_calls": []},
        })

    mode_result = await _try_question_mode_route(question_mode, text_value, history, session_id, db)
    if mode_result:
        log_turn(
            session_id=session_id,
            user_text=text_value,
            intent=str((mode_result.get("debug") or {}).get("intent") or "general"),
            provider="question_mode",
            route=str((mode_result.get("debug") or {}).get("route") or "mode"),
            tool_calls=(mode_result.get("debug") or {}).get("tool_calls"),
        )
        return _finish(mode_result)

    fast_price = _national_price_fast_tool(text_value, history, session_id, question_mode)
    if fast_price:
        name, args = fast_price
        result = await _run_national_price_tool(db, name, args)
        reply, data_card, report = _reply_from_national_tool(name, result)
        _remember_turn(session_id, text_value, reply)
        log_turn(
            session_id=session_id,
            user_text=text_value,
            intent="national_price",
            provider="fast_path",
            route="price_fast",
            tool_calls=[{"name": name, "arguments": args}],
        )
        return _finish({
            "reply": reply,
            "data_card": data_card,
            "report_content": report,
            "session_id": session_id or f"chat-{datetime.utcnow().strftime('%Y%m%d')}",
            "debug": {
                "provider": "fast_path",
                "route": "price_fast",
                "intent": "national_price",
                "tool_calls": [{"name": name, "arguments": args, "result_preview": _preview(result)}],
            },
        })

    if _VAGUE_ORDER_PATTERN.match((text_value or "").strip()):
        d = _today().isoformat()
        name, args = "search_orders", {"start_date": d, "end_date": d, "limit": 20}
        result = await _dispatch_tool(name, args, db)
        reply = (
            f"未指定时间范围，已按**今天（{d}）**查询订单列表。"
            f"{result.get('summary') or ''}"
        ).strip()
        _remember_turn(session_id, text_value, reply)
        log_turn(
            session_id=session_id,
            user_text=text_value,
            intent="analytics",
            provider="local_rule",
            route="vague_orders_default",
            tool_calls=[{"name": name, "arguments": args}],
        )
        data_card, report = _artifact_from_tool(result)
        return _finish({
            "reply": reply,
            "data_card": data_card,
            "report_content": report,
            "session_id": session_id or f"chat-{datetime.utcnow().strftime('%Y%m%d')}",
            "debug": {
                "provider": "local_rule",
                "route": "vague_orders_default",
                "intent": "analytics",
                "tool_calls": [{"name": name, "arguments": args, "result_preview": _preview(result)}],
            },
        })

    chat_intent = classify_chat_intent(text_value)
    seed_memory_from_history(session_id, history)
    entities = apply_memory_entities(load_session(session_id), _extract_entities_from_history(history, text_value))

    if not (settings.ai_api_key or "").strip():
        result = await _local_answer(text_value, db, entities)
        result["session_id"] = session_id or result.get("session_id")
        _remember_turn(session_id, text_value, result.get("reply") or "")
        finished = _finish(result)
        log_turn(
            session_id=session_id,
            user_text=text_value,
            intent=chat_intent,
            provider="local_rule_fallback",
            route="fallback",
            fallback=True,
            tool_calls=(result.get("debug") or {}).get("tool_calls"),
            duration_ms=(finished.get("debug") or {}).get("duration_ms"),
        )
        return finished

    try:
        out = await _run_llm_agent(text_value, history, db, session_id, chat_intent)
        finished = _finish(out)
        log_turn(
            session_id=session_id,
            user_text=text_value,
            intent=chat_intent,
            provider="dashscope",
            route="llm_first",
            tool_calls=(out.get("debug") or {}).get("tool_calls"),
            synthesis=bool((out.get("debug") or {}).get("synthesis")),
            duration_ms=(finished.get("debug") or {}).get("duration_ms"),
        )
        return finished
    except Exception as exc:  # noqa: BLE001
        result = await _local_answer(text_value, db, entities)
        message = str(exc) or repr(exc) or exc.__class__.__name__
        dbg = result.setdefault("debug", {})
        dbg["provider_error"] = message[:220]
        dbg["provider"] = "local_rule_fallback"
        dbg["route"] = "fallback"
        dbg["intent"] = chat_intent
        result["session_id"] = session_id or result.get("session_id")
        _remember_turn(session_id, text_value, result.get("reply") or "")
        finished = _finish(result)
        log_turn(
            session_id=session_id,
            user_text=text_value,
            intent=chat_intent,
            provider="local_rule_fallback",
            route="fallback",
            fallback=True,
            error=message,
            tool_calls=dbg.get("tool_calls"),
            duration_ms=(finished.get("debug") or {}).get("duration_ms"),
        )
        return finished


def _national_price_fast_tool(
    text_value: str,
    history: list[ChatMessage],
    session_id: Optional[str],
    question_mode: Optional[str],
) -> Optional[tuple[str, dict[str, Any]]]:
    if (question_mode or "auto").strip().lower() == "system_data":
        return None
    return (
        _fast_price_query_from_text(text_value)
        or _fast_price_followup(text_value, history, session_id)
        or _fast_historical_price_from_text(text_value)
    )


async def _iter_forecast_training_sse(
    db: AsyncSession,
    name: str,
    args: dict[str, Any],
):
    if name not in ("get_national_ag_forecast_price", "get_xinfadi_forecast_price"):
        return
    raw = str(args.get("query_text") or "").strip()
    parsed = national_price_svc.parse_price_query(raw)
    if args.get("product_query"):
        parsed = national_price_svc.ParsedPriceQuery(
            raw_text=raw,
            product_query=str(args.get("product_query")),
            goods_name_hint=parsed.goods_name_hint,
            grade_hint=parsed.grade_hint,
            target_date=parsed.target_date,
            intent=parsed.intent,
        )
    sku_key, _, _ = await national_price_svc.resolve_sku_from_parsed(db, parsed)
    if not sku_key or await zg_train_svc.snapshot_exists(db, sku_key):
        return
    async for ev in zg_train_svc.iter_forecast_train_progress(db, sku_key):
        yield sse_training_phase(ev)


async def _finish_national_price_fast(
    db: AsyncSession,
    text_value: str,
    session_id: Optional[str],
    name: str,
    args: dict[str, Any],
    *,
    route: str = "price_fast",
    provider: str = "fast_path",
) -> dict[str, Any]:
    result = await _run_national_price_tool(db, name, args)
    reply, data_card, report = _reply_from_national_tool(name, result)
    _remember_turn(session_id, text_value, reply)
    return {
        "reply": reply,
        "data_card": data_card,
        "report_content": report,
        "session_id": session_id or f"chat-{datetime.utcnow().strftime('%Y%m%d')}",
        "debug": {
            "provider": provider,
            "route": route,
            "intent": "national_price",
            "tool_calls": [{"name": name, "arguments": args, "result_preview": _preview(result)}],
        },
    }


@router.post("")
async def chat(
    body: ChatRequest,
    _=Depends(_monitor_guard),
    db: AsyncSession = Depends(get_db),
):
    text_value = _latest_user_text(body)
    return await _dispatch(text_value, body.messages, db, body.session_id, body.question_mode)


@router.post("/stream")
async def chat_stream(
    body: ChatRequest,
    _=Depends(_monitor_guard),
    db: AsyncSession = Depends(get_db),
):
    text_value = _latest_user_text(body)
    history = body.messages
    session_id = body.session_id
    question_mode = body.question_mode

    async def events():
        import json as _json

        mode = (question_mode or "auto").strip().lower()
        if mode == "national_price":
            tool_pair = (
                _fast_price_query_from_text(text_value)
                or _fast_price_followup(text_value, history, session_id)
                or _fast_historical_price_from_text(text_value)
                or _national_price_tool_from_text(text_value)
            )
            if tool_pair:
                name, args = tool_pair
                yield f"event: phase\ndata: {_json.dumps({'phase': 'running', 'message': '正在查询全国农产品价格…'}, ensure_ascii=False)}\n\n"
                async for line in _iter_forecast_training_sse(db, name, args):
                    yield line
                result = await _finish_national_price_fast(
                    db, text_value, session_id, name, args, route="mode_national_price", provider="question_mode"
                )
                for line in sse_events_from_result(result):
                    yield line
                return

        fast = _national_price_fast_tool(text_value, history, session_id, question_mode)
        if fast:
            name, args = fast
            yield f"event: phase\ndata: {_json.dumps({'phase': 'running', 'message': '正在查询全国农产品价格…'}, ensure_ascii=False)}\n\n"
            async for line in _iter_forecast_training_sse(db, name, args):
                yield line
            result = await _finish_national_price_fast(db, text_value, session_id, name, args)
            log_turn(
                session_id=session_id,
                user_text=text_value,
                intent="national_price",
                provider="fast_path",
                route="price_fast",
                tool_calls=(result.get("debug") or {}).get("tool_calls"),
            )
            for line in sse_events_from_result(result):
                yield line
            return

        result = await _dispatch(text_value, history, db, session_id, question_mode)
        for line in sse_events_from_result(result):
            yield line

    return StreamingResponse(events(), media_type="text/event-stream")


@router.get("/catalog")
async def catalog(_=Depends(_monitor_guard), db: AsyncSession = Depends(get_db)):
    schema = await _tool_schema_overview(db, {})
    sample = await _tool_kpi_summary(db, {"scope": "today"})
    return {
        "updated_at": datetime.utcnow().isoformat(),
        "provider_configured": bool((settings.ai_api_key or "").strip()),
        "tables": schema.get("rows", []),
        "api_samples": {"kpi": sample},
    }


@router.post("/catalog/refresh")
async def catalog_refresh(_=Depends(_monitor_guard), db: AsyncSession = Depends(get_db)):
    return await catalog(_=None, db=db)


@router.get("/session/{session_id}/debug")
async def session_debug(session_id: str, _=Depends(_monitor_guard)):
    mem = load_session(session_id)
    return {
        "session_id": session_id,
        "session_memory": mem.to_debug(),
        "history_turns": len(_SESSION_HISTORY.get(session_id) or []),
    }


@router.get("/docs/search")
async def docs_search(
    q: str = Query("", min_length=1),
    limit: int = Query(5, ge=1, le=20),
    _=Depends(_monitor_guard),
):
    hits = search_docs(q, top_k=limit, audience="monitor")
    return {"query": q, "total": len(hits), "items": hits, "citations": citations_from_rag(hits)}


def _content_disposition(filename: str, ext: str) -> str:
    """中文文件名安全编码（HTTP 头是 latin-1，必须用 RFC 5987 filename*）。"""
    from urllib.parse import quote

    name = f"{(filename or '导出').strip() or '导出'}.{ext}"
    ascii_fallback = re.sub(r"[^A-Za-z0-9._-]", "_", name) or f"export.{ext}"
    return f"attachment; filename=\"{ascii_fallback}\"; filename*=UTF-8''{quote(name)}"


@router.post("/report/export")
async def export_report(
    body: ExportRequest,
    _=Depends(_monitor_guard),
):
    title = body.title.strip() or "监管分析报告"
    content = (body.markdown or body.content or "").strip() or "# 监管分析报告\n\n暂无报告内容。"
    fmt = (body.format or "docx").lower()
    if fmt in ("xlsx", "excel"):
        try:
            import xlsxwriter

            buf = io.BytesIO()
            wb = xlsxwriter.Workbook(buf, {"in_memory": True})
            ws = wb.add_worksheet("数据")
            cols = body.columns or ([{"key": k, "label": k} for k in (body.rows[0].keys() if body.rows else [])])
            head_fmt = wb.add_format({"bold": True, "bg_color": "#0a2a33", "font_color": "#9cf0ff", "border": 1})
            for ci, col in enumerate(cols):
                ws.write(0, ci, str(col.get("label") or col.get("key") or ""), head_fmt)
                ws.set_column(ci, ci, 18)
            for ri, row in enumerate(body.rows, start=1):
                for ci, col in enumerate(cols):
                    val = row.get(col.get("key"))
                    ws.write(ri, ci, "" if val is None else (val if isinstance(val, (int, float)) else str(val)))
            wb.close()
            filename = body.filename or title
            return Response(
                content=buf.getvalue(),
                media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                headers={"Content-Disposition": _content_disposition(filename, "xlsx")},
            )
        except Exception:
            fmt = "md"
    if fmt == "docx":
        try:
            from docx import Document

            doc = Document()
            doc.add_heading(title, 0)
            for line in content.splitlines():
                clean = line.strip()
                if not clean:
                    continue
                if clean.startswith("# "):
                    doc.add_heading(clean[2:].strip(), level=1)
                elif clean.startswith("## "):
                    doc.add_heading(clean[3:].strip(), level=2)
                elif clean.startswith("- "):
                    doc.add_paragraph(clean[2:].strip(), style="List Bullet")
                else:
                    doc.add_paragraph(clean)
            buf = io.BytesIO()
            doc.save(buf)
            filename = body.filename or title
            return Response(
                content=buf.getvalue(),
                media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                headers={"Content-Disposition": _content_disposition(filename, "docx")},
            )
        except Exception:
            fmt = "md"
    if fmt == "pptx":
        try:
            from pptx import Presentation

            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = title
            slide.placeholders[1].text = "\n".join([ln.strip("# ").strip() for ln in content.splitlines() if ln.strip()][:8])
            buf = io.BytesIO()
            prs.save(buf)
            filename = body.filename or title
            return Response(
                content=buf.getvalue(),
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                headers={"Content-Disposition": _content_disposition(filename, "pptx")},
            )
        except Exception:
            fmt = "md"
    payload = f"{title}\n\n{content}".encode("utf-8")
    filename = body.filename or title
    return Response(
        content=payload,
        media_type="text/markdown; charset=utf-8",
        headers={"Content-Disposition": _content_disposition(filename, "md")},
    )
